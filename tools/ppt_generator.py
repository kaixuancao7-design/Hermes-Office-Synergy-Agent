"""PPT生成器 - 基于python-pptx的PPT生成工具"""

from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import uuid
from src.logging_config import get_logger

logger = get_logger("ppt.generator")


class PPTGeneratorBase:
    """PPT生成器基类"""
    
    def __init__(self):
        self.output_dir = "output/ppt"
        os.makedirs(self.output_dir, exist_ok=True)
    
    def generate_ppt(self, title: str, slides: List[Dict[str, Any]], output_path: Optional[str] = None) -> str:
        """
        生成PPT文件
        
        Args:
            title: PPT标题
            slides: 幻灯片列表
            output_path: 输出路径（可选）
        
        Returns:
            生成的PPT文件路径
        """
        try:
            prs = Presentation()
            
            self._add_title_slide(prs, title)
            
            for slide_data in slides:
                slide_type = slide_data.get("type", "content")
                if slide_type == "title":
                    self._add_title_slide(prs, slide_data.get("title", ""), slide_data.get("subtitle", ""))
                elif slide_type == "content":
                    self._add_content_slide(prs, slide_data)
                elif slide_type == "list":
                    self._add_list_slide(prs, slide_data)
                elif slide_type == "image":
                    self._add_image_slide(prs, slide_data)
                else:
                    self._add_content_slide(prs, slide_data)
            
            if output_path:
                output_file = output_path
            else:
                file_name = f"{title}_{uuid.uuid4().hex[:8]}.pptx"
                output_file = os.path.join(self.output_dir, file_name)
            
            prs.save(output_file)
            logger.info(f"PPT生成成功: {output_file}")
            
            return output_file
        
        except Exception as e:
            logger.error(f"PPT生成失败: {str(e)}")
            raise
    
    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str = ""):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        subtitle_placeholder.text = subtitle
        
        title_text_frame = title_placeholder.text_frame
        for paragraph in title_text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 51, 102)
    
    def _add_content_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        
        slide.shapes.title.text = title
        
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content
    
    def _add_list_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide_data.get("title", "")
        items = slide_data.get("items", [])
        
        slide.shapes.title.text = title
        
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
    
    def _add_image_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """添加图片幻灯片"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide_data.get("title", "")
        image_path = slide_data.get("image_path", "")
        
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_text = title_box.text_frame
        title_paragraph = title_text.add_paragraph()
        title_paragraph.text = title
        title_paragraph.font.size = Pt(24)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
    
    def generate_from_outline(self, title: str, outline: List[Dict[str, Any]]) -> str:
        """
        从大纲生成PPT
        
        Args:
            title: PPT标题
            outline: 大纲结构列表
        
        Returns:
            生成的文件路径
        """
        slides = self._convert_outline_to_slides(title, outline)
        return self.generate_ppt(title, slides)
    
    def _convert_outline_to_slides(self, title: str, outline: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        将大纲结构转换为幻灯片列表
        
        Args:
            title: PPT标题
            outline: 大纲结构
        
        Returns:
            幻灯片列表
        """
        slides = []
        
        slides.append({
            "type": "title",
            "title": title,
            "subtitle": ""
        })
        
        for item in outline:
            slide_title = item.get("title", "")
            content = item.get("content", "")
            items = item.get("items", [])
            
            if items:
                slides.append({
                    "type": "list",
                    "title": slide_title,
                    "items": items
                })
            else:
                slides.append({
                    "type": "content",
                    "title": slide_title,
                    "content": content
                })
        
        slides.append({
            "type": "content",
            "title": "谢谢",
            "content": ""
        })
        
        return slides


class GeneratePPT(PPTGeneratorBase):
    """PPT生成器（向后兼容）"""
    pass


# 扩展功能：支持从大纲生成PPT
def generate_ppt_from_outline(title: str, outline: List[Dict[str, Any]], output_path: Optional[str] = None) -> str:
    """
    从大纲生成PPT（便捷函数）
    
    Args:
        title: PPT标题
        outline: 大纲结构
        output_path: 输出路径
    
    Returns:
        生成的文件路径
    """
    generator = PPTGeneratorBase()
    return generator.generate_ppt(title, outline, output_path)