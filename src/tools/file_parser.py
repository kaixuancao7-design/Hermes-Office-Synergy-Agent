"""文件解析工具 - 支持多种文件格式的内容提取"""
from typing import Dict, Any, Optional, List
from src.logging_config import get_logger

logger = get_logger("tool")


class FileParser:
    """文件解析器 - 支持多种文件格式
    
    支持的文件类型：
    - .txt, .md: 纯文本文件
    - .docx: Word文档
    - .xlsx, .xls: Excel表格
    - .pptx: PowerPoint演示文稿
    - .pdf: PDF文档
    - .json: JSON文件
    - .xml: XML文件
    - .html: HTML文件
    - .rtf: RTF文档
    - .csv, .tsv: 表格文件
    """
    
    def __init__(self):
        self.parsers = {
            ".txt": self._parse_text,
            ".md": self._parse_markdown,
            ".docx": self._parse_docx,
            ".xlsx": self._parse_xlsx,
            ".xls": self._parse_xls,
            ".pptx": self._parse_pptx,
            ".pdf": self._parse_pdf,
            ".json": self._parse_json,
            ".xml": self._parse_xml,
            ".html": self._parse_html,
            ".rtf": self._parse_rtf,
            ".csv": self._parse_csv,
            ".tsv": self._parse_tsv,
        }
    
    def parse(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """
        解析文件内容
        
        Args:
            file_path: 文件路径（用于判断扩展名）
            file_content: 文件二进制内容（可选，优先使用）
            
        Returns:
            解析结果字典，包含内容、元数据等
        """
        import os
        
        # 获取文件扩展名
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext not in self.parsers:
            return {
                "success": False,
                "error": f"不支持的文件类型: {ext}",
                "content": "",
                "metadata": {}
            }
        
        try:
            return self.parsers[ext](file_path, file_content)
        except Exception as e:
            logger.error(f"解析文件 {file_path} 失败: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "content": "",
                "metadata": {}
            }
    
    def _parse_text(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析纯文本文件"""
        content = self._read_file_content(file_path, file_content)
        return {
            "success": True,
            "content": content,
            "metadata": {"type": "text", "encoding": "utf-8"}
        }
    
    def _parse_markdown(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析Markdown文件"""
        content = self._read_file_content(file_path, file_content)
        return {
            "success": True,
            "content": content,
            "metadata": {"type": "markdown", "encoding": "utf-8"}
        }
    
    def _parse_docx(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析Word文档"""
        try:
            from docx import Document
            from io import BytesIO
            
            if file_content:
                doc = Document(BytesIO(file_content))
            else:
                doc = Document(file_path)
            
            content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # 获取文档属性
            metadata = {
                "type": "docx",
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "created": doc.core_properties.created.isoformat() if doc.core_properties.created else "",
                "modified": doc.core_properties.modified.isoformat() if doc.core_properties.modified else ""
            }
            
            return {
                "success": True,
                "content": content,
                "metadata": metadata
            }
        except ImportError:
            logger.warning("python-docx 未安装，使用纯文本解析")
            return self._parse_text(file_path, file_content)
        except Exception as e:
            logger.error(f"解析docx失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def _parse_xlsx(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析Excel文件(.xlsx)"""
        try:
            import pandas as pd
            from io import BytesIO
            
            if file_content:
                df = pd.read_excel(BytesIO(file_content), engine="openpyxl")
            else:
                df = pd.read_excel(file_path, engine="openpyxl")
            
            content = self._convert_dataframe_to_text(df)
            sheet_names = self._get_excel_sheet_names(file_path, file_content)
            
            return {
                "success": True,
                "content": content,
                "metadata": {"type": "xlsx", "sheet_names": sheet_names, "rows": len(df), "columns": len(df.columns)}
            }
        except ImportError:
            logger.warning("pandas/openpyxl 未安装")
            return {"success": False, "error": "缺少依赖库", "content": "", "metadata": {}}
        except Exception as e:
            logger.error(f"解析xlsx失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def _parse_xls(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析Excel文件(.xls)"""
        try:
            import pandas as pd
            from io import BytesIO
            
            if file_content:
                df = pd.read_excel(BytesIO(file_content), engine="xlrd")
            else:
                df = pd.read_excel(file_path, engine="xlrd")
            
            content = self._convert_dataframe_to_text(df)
            
            return {
                "success": True,
                "content": content,
                "metadata": {"type": "xls", "rows": len(df), "columns": len(df.columns)}
            }
        except ImportError:
            logger.warning("pandas/xlrd 未安装")
            return {"success": False, "error": "缺少依赖库", "content": "", "metadata": {}}
        except Exception as e:
            logger.error(f"解析xls失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def _parse_pptx(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析PowerPoint文件"""
        try:
            from pptx import Presentation
            from io import BytesIO
            
            if file_content:
                prs = Presentation(BytesIO(file_content))
            else:
                prs = Presentation(file_path)
            
            slides_content = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slides_content.append("\n".join(slide_text))
            
            content = "\n\n---\n\n".join(slides_content)
            
            return {
                "success": True,
                "content": content,
                "metadata": {"type": "pptx", "slides": len(prs.slides)}
            }
        except ImportError:
            logger.warning("python-pptx 未安装")
            return {"success": False, "error": "缺少依赖库", "content": "", "metadata": {}}
        except Exception as e:
            logger.error(f"解析pptx失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def _parse_pdf(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析PDF文件"""
        try:
            from PyPDF2 import PdfReader
            from io import BytesIO
            
            if file_content:
                reader = PdfReader(BytesIO(file_content))
            else:
                reader = PdfReader(file_path)
            
            content = "\n".join([page.extract_text() for page in reader.pages])
            
            return {
                "success": True,
                "content": content,
                "metadata": {"type": "pdf", "pages": len(reader.pages)}
            }
        except ImportError:
            logger.warning("PyPDF2 未安装")
            return {"success": False, "error": "缺少依赖库", "content": "", "metadata": {}}
        except Exception as e:
            logger.error(f"解析pdf失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def _parse_json(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析JSON文件"""
        content = self._read_file_content(file_path, file_content)
        try:
            import json
            data = json.loads(content)
            return {
                "success": True,
                "content": json.dumps(data, ensure_ascii=False, indent=2),
                "metadata": {"type": "json", "parsed": data}
            }
        except json.JSONDecodeError as e:
            return {"success": False, "error": f"JSON解析错误: {str(e)}", "content": content, "metadata": {"type": "json"}}
    
    def _parse_xml(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析XML文件"""
        content = self._read_file_content(file_path, file_content)
        try:
            import xml.etree.ElementTree as ET
            tree = ET.ElementTree(ET.fromstring(content))
            root = tree.getroot()
            
            # 提取标签信息
            tags = []
            for elem in root.iter():
                tags.append(elem.tag)
            
            return {
                "success": True,
                "content": content,
                "metadata": {"type": "xml", "root_tag": root.tag, "tags_count": len(set(tags))}
            }
        except Exception as e:
            return {"success": False, "error": f"XML解析错误: {str(e)}", "content": content, "metadata": {"type": "xml"}}
    
    def _parse_html(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析HTML文件"""
        content = self._read_file_content(file_path, file_content)
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(content, "html.parser")
            text_content = soup.get_text(strip=True, separator="\n")
            
            return {
                "success": True,
                "content": text_content,
                "metadata": {"type": "html", "title": soup.title.string if soup.title else ""}
            }
        except ImportError:
            logger.warning("beautifulsoup4 未安装")
            return {"success": True, "content": content, "metadata": {"type": "html"}}
        except Exception as e:
            logger.error(f"解析html失败: {str(e)}")
            return {"success": False, "error": str(e), "content": content, "metadata": {"type": "html"}}
    
    def _parse_rtf(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析RTF文件"""
        content = self._read_file_content(file_path, file_content)
        try:
            from striprtf.striprtf import rtf_to_text
            text_content = rtf_to_text(content)
            
            return {
                "success": True,
                "content": text_content,
                "metadata": {"type": "rtf"}
            }
        except ImportError:
            logger.warning("striprtf 未安装")
            return {"success": False, "error": "缺少依赖库", "content": "", "metadata": {"type": "rtf"}}
        except Exception as e:
            logger.error(f"解析rtf失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {"type": "rtf"}}
    
    def _parse_csv(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析CSV文件"""
        content = self._read_file_content(file_path, file_content)
        try:
            import csv
            from io import StringIO
            
            reader = csv.reader(StringIO(content))
            rows = list(reader)
            
            return {
                "success": True,
                "content": content,
                "metadata": {"type": "csv", "rows": len(rows), "columns": len(rows[0]) if rows else 0}
            }
        except Exception as e:
            logger.error(f"解析csv失败: {str(e)}")
            return {"success": False, "error": str(e), "content": content, "metadata": {"type": "csv"}}
    
    def _parse_tsv(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        """解析TSV文件"""
        content = self._read_file_content(file_path, file_content)
        try:
            import csv
            from io import StringIO
            
            reader = csv.reader(StringIO(content), delimiter="\t")
            rows = list(reader)
            
            return {
                "success": True,
                "content": content,
                "metadata": {"type": "tsv", "rows": len(rows), "columns": len(rows[0]) if rows else 0}
            }
        except Exception as e:
            logger.error(f"解析tsv失败: {str(e)}")
            return {"success": False, "error": str(e), "content": content, "metadata": {"type": "tsv"}}
    
    def _read_file_content(self, file_path: str, file_content: Optional[bytes] = None) -> str:
        """读取文件内容为字符串"""
        if file_content:
            try:
                return file_content.decode("utf-8")
            except UnicodeDecodeError:
                return file_content.decode("gbk", errors="ignore")
        else:
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()
            except UnicodeDecodeError:
                with open(file_path, "r", encoding="gbk") as f:
                    return f.read()
    
    def _convert_dataframe_to_text(self, df) -> str:
        """将DataFrame转换为文本格式"""
        lines = []
        # 添加列名
        lines.append("\t".join(str(col) for col in df.columns))
        # 添加数据行
        for _, row in df.iterrows():
            lines.append("\t".join(str(val) for val in row))
        return "\n".join(lines)
    
    def _get_excel_sheet_names(self, file_path: str, file_content: Optional[bytes] = None) -> List[str]:
        """获取Excel文件的所有工作表名称"""
        try:
            from openpyxl import load_workbook
            from io import BytesIO
            
            if file_content:
                wb = load_workbook(filename=BytesIO(file_content), read_only=True)
            else:
                wb = load_workbook(filename=file_path, read_only=True)
            
            return wb.sheetnames
        except Exception:
            return []
    
    def get_supported_formats(self) -> List[str]:
        """获取支持的文件格式列表"""
        return list(self.parsers.keys())


# 单例实例
file_parser = FileParser()


def parse_file(file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
    """
    解析文件内容（便捷函数）
    
    Args:
        file_path: 文件路径
        file_content: 文件二进制内容（可选）
    
    Returns:
        解析结果
    """
    return file_parser.parse(file_path, file_content)


def extract_text_from_file(file_path: str, file_content: Optional[bytes] = None) -> str:
    """
    从文件中提取文本内容（便捷函数）
    
    Args:
        file_path: 文件路径
        file_content: 文件二进制内容（可选）
    
    Returns:
        提取的文本内容，如果失败返回空字符串
    """
    result = file_parser.parse(file_path, file_content)
    return result.get("content", "") if result.get("success") else ""
