from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from src.logging_config import get_logger
from src.utils import generate_id, get_timestamp
from src.tools.base import BaseTool, ToolSchema
from src.tools.registry import register_tool
from pydantic import Field

logger = get_logger("tool.ppt")


class GeneratePPTSchema(ToolSchema):
    """生成PPT工具参数Schema"""
    title: str = Field(description="PPT标题", default="Untitled Presentation")
    slides: List[Dict[str, Any]] = Field(description="幻灯片列表", default_factory=list)


class GeneratePPTFromOutlineSchema(ToolSchema):
    """从大纲生成PPT工具参数Schema"""
    title: str = Field(description="PPT标题", default="Untitled Presentation")
    outline: List[Dict[str, Any]] = Field(description="大纲结构", default_factory=list)


class PPTGeneratorBase:
    """PPT生成器核心实现 - 基于python-pptx库"""

    def __init__(self):
        self.output_dir = "output/ppt"
        import os
        os.makedirs(self.output_dir, exist_ok=True)

    def _parse_color(self, color_str: str) -> RGBColor:
        """解析颜色字符串为RGBColor"""
        if isinstance(color_str, RGBColor):
            return color_str
        try:
            if color_str.startswith("#"):
                color_str = color_str[1:]
            r = int(color_str[0:2], 16)
            g = int(color_str[2:4], 16)
            b = int(color_str[4:6], 16)
            return RGBColor(r, g, b)
        except:
            return RGBColor(0, 51, 102)
    
    def generate_ppt(self, title: str, slides: List[Dict[str, Any]], 
                     output_path: Optional[str] = None) -> str:
        """
        生成PPT文件
        
        Args:
            title: PPT标题
            slides: 幻灯片列表，每个元素包含:
                - type: 幻灯片类型 (title, content, image, chart, blank)
                - title: 幻灯片标题
                - content: 内容（文本或列表）
                - layout: 布局类型
            output_path: 输出路径，默认为自动生成
        
        Returns:
            生成的PPT文件路径
        """
        try:
            # 创建演示文稿
            prs = Presentation()
            
            # 添加封面页
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]
            
            title_placeholder.text = title
            subtitle_placeholder.text = f"生成时间: {get_timestamp()}"
            
            # 设置标题样式
            title_text_frame = title_placeholder.text_frame
            for paragraph in title_text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(36)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 51, 102)
            
            # 添加内容页
            for slide_data in slides:
                self._add_slide(prs, slide_data)
            
            # 生成输出路径
            if not output_path:
                output_path = f"{self.output_dir}/{generate_id()[:8]}_{title[:20]}.pptx"
            
            # 确保目录存在
            import os
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            # 保存文件
            prs.save(output_path)
            logger.info(f"PPT generated successfully: {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"PPT generation failed: {str(e)}", exc_info=True)
            raise
    
    def _add_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """添加单个幻灯片"""
        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        style_meta = slide_data.get("_style", {})

        color_scheme = style_meta.get("color_scheme", {})
        font_family = style_meta.get("font_family", {})

        primary_color = self._parse_color(color_scheme.get("primary", "003366"))
        font_heading = font_family.get("heading", "Microsoft YaHei")
        font_body = font_family.get("body", "Microsoft YaHei")

        if slide_type == "title":
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            if isinstance(content, str):
                slide.placeholders[1].text = content

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = primary_color
                            run.font.name = font_heading

        elif slide_type == "content":
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title

            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame

            if isinstance(content, list):
                for i, item in enumerate(content):
                    paragraph = text_frame.add_paragraph()
                    if i > 0:
                        paragraph.level = 1
                    paragraph.text = str(item)
            else:
                text_frame.text = str(content)

            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = font_body
        
        elif slide_type == "bullet":
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title

            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame

            if isinstance(content, list):
                for item in content:
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = str(item)
                    paragraph.level = 0

            for paragraph in text_frame.paragraphs:
                paragraph.font.name = font_body

        elif slide_type == "image":
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
            
            if content and isinstance(content, str) and content.startswith("http"):
                try:
                    import requests
                    from io import BytesIO
                    
                    response = requests.get(content)
                    image_stream = BytesIO(response.content)
                    left = top = Inches(1)
                    height = Inches(5)
                    slide.shapes.add_picture(image_stream, left, top, height=height)
                except Exception as e:
                    logger.warning(f"Failed to add image: {e}")
        
        elif slide_type == "chart":
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
                if title_shape.has_text_frame:
                    for paragraph in title_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = primary_color
                            run.font.name = font_heading

            if isinstance(content, list) and len(content) > 0:
                rows = len(content) + 1
                cols = len(content[0]) if isinstance(content[0], list) else 2

                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                height = Inches(4)

                table = slide.shapes.add_table(rows, cols, left, top, width, height).table

                header_row = table.rows[0]
                for i in range(cols):
                    header_row.cells[i].text = f"列{i+1}"
                    header_row.cells[i].fill.solid()
                    header_row.cells[i].fill.fore_color.rgb = primary_color
                    for paragraph in header_row.cells[i].text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                
                for i, row_data in enumerate(content):
                    row = table.rows[i+1]
                    if isinstance(row_data, list):
                        for j, cell_data in enumerate(row_data[:cols]):
                            row.cells[j].text = str(cell_data)
                    else:
                        row.cells[0].text = str(row_data)
        
        elif slide_type == "blank":
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            if title:
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = title
    
    def generate_from_outline(self, title: str, outline: List[Dict[str, Any]]) -> str:
        """
        从大纲生成PPT
        
        Args:
            title: PPT标题
            outline: 大纲结构，每个元素包含:
                - title: 章节标题
                - content: 章节内容（文本或子列表）
        
        Returns:
            生成的PPT文件路径
        """
        slides = []
        
        for section in outline:
            section_title = section.get("title", "")
            section_content = section.get("content", "")
            
            slides.append({
                "type": "title",
                "title": section_title
            })
            
            if section_content:
                if isinstance(section_content, list):
                    slides.append({
                        "type": "bullet",
                        "title": section_title,
                        "content": section_content
                    })
                else:
                    slides.append({
                        "type": "content",
                        "title": section_title,
                        "content": section_content
                    })
        
        return self.generate_ppt(title, slides)


@register_tool("generate_ppt")
class GeneratePPT(BaseTool):
    """生成PPT工具"""
    
    description = "根据幻灯片列表生成PPT文件"
    schema = GeneratePPTSchema
    
    def __init__(self, executor=None):
        self.generator = PPTGeneratorBase()
        self.executor = executor
    
    def execute(self, params: Dict[str, Any]) -> Dict[str, Any]:
        title = params.get("title", "Untitled Presentation")
        slides = params.get("slides", [])
        
        if not slides or not isinstance(slides, list):
            return {"success": False, "error": "参数错误：slides必须是非空数组"}
        
        try:
            output_path = self.generator.generate_ppt(title, slides)
            logger.info(f"PPT generated successfully: {output_path}")
            return {
                "success": True,
                "result": {
                    "file_path": output_path,
                    "title": title,
                    "slides_count": len(slides)
                }
            }
        except Exception as e:
            logger.error(f"PPT generation failed: {str(e)}", exc_info=True)
            return {"success": False, "error": f"PPT生成失败: {str(e)}"}


@register_tool("generate_ppt_from_outline")
class GeneratePPTFromOutline(BaseTool):
    """从大纲生成PPT工具"""
    
    description = "根据大纲结构生成PPT文件"
    schema = GeneratePPTFromOutlineSchema
    
    def __init__(self, executor=None):
        self.generator = PPTGeneratorBase()
        self.executor = executor
    
    def execute(self, params: Dict[str, Any]) -> Dict[str, Any]:
        title = params.get("title", "Untitled Presentation")
        outline = params.get("outline", [])
        
        if not outline or not isinstance(outline, list):
            return {"success": False, "error": "参数错误：outline必须是非空数组"}
        
        try:
            output_path = self.generator.generate_from_outline(title, outline)
            logger.info(f"PPT generated from outline: {output_path}")
            return {
                "success": True,
                "result": {
                    "file_path": output_path,
                    "title": title,
                    "chapters_count": len(outline)
                }
            }
        except Exception as e:
            logger.error(f"PPT generation from outline failed: {str(e)}", exc_info=True)
            return {"success": False, "error": f"PPT生成失败: {str(e)}"}
