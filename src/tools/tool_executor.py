from typing import Dict, Any, Optional, List
from src.infrastructure.sandbox import sandbox
from src.utils import generate_id, get_timestamp
from src.logging_config import get_logger
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

logger = get_logger("tool")


class PPTGenerator:
    """PPT生成器 - 基于python-pptx库"""
    
    def __init__(self):
        self.output_dir = "output/ppt"
        import os
        os.makedirs(self.output_dir, exist_ok=True)
    
    def generate_ppt(self, title: str, slides: List[Dict[str, Any]], output_path: Optional[str] = None) -> str:
        """
        生成PPT文件
        
        Args:
            title: PPT标题
            slides: 幻灯片列表，每个元素包含:
                - type: 幻灯片类型 (title, content, image, chart, blank)
                - title: 幻灯片标题
                - content: 内容（文本或列表）
                - layout: 布局类型
            output_path: 输出路径，默认为自动生成
        
        Returns:
            生成的PPT文件路径
        """
        try:
            # 创建演示文稿
            prs = Presentation()
            
            # 添加封面页
            slide_layout = prs.slide_layouts[0]  # 标题页布局
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]
            
            title_placeholder.text = title
            subtitle_placeholder.text = f"生成时间: {get_timestamp()}"
            
            # 设置标题样式
            title_text_frame = title_placeholder.text_frame
            for paragraph in title_text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(36)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 51, 102)
            
            # 添加内容页
            for slide_data in slides:
                self._add_slide(prs, slide_data)
            
            # 生成输出路径
            if not output_path:
                output_path = f"{self.output_dir}/{generate_id()[:8]}_{title[:20]}.pptx"
            
            # 确保目录存在
            import os
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            # 保存文件
            prs.save(output_path)
            logger.info(f"PPT generated successfully: {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"PPT generation failed: {str(e)}", exc_info=True)
            raise
    
    def _add_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """添加单个幻灯片"""
        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        
        if slide_type == "title":
            # 标题页
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            if isinstance(content, str):
                slide.placeholders[1].text = content
        
        elif slide_type == "content":
            # 内容页（标题+正文）
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            
            if isinstance(content, list):
                # 列表内容
                for i, item in enumerate(content):
                    if i == 0:
                        paragraph = text_frame.add_paragraph()
                    else:
                        paragraph = text_frame.add_paragraph()
                        paragraph.level = 1
                    paragraph.text = str(item)
            else:
                # 文本内容
                text_frame.text = str(content)
            
            # 设置字体
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(18)
        
        elif slide_type == "bullet":
            # 项目符号页
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            
            if isinstance(content, list):
                for i, item in enumerate(content):
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = str(item)
                    paragraph.level = 0
        
        elif slide_type == "image":
            # 图片页
            slide_layout = prs.slide_layouts[5]  # 空白页
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加标题
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
            
            # 如果有图片路径，添加图片
            if content and isinstance(content, str) and content.startswith("http"):
                # 下载图片并添加
                try:
                    import requests
                    from io import BytesIO
                    
                    response = requests.get(content)
                    image_stream = BytesIO(response.content)
                    left = top = Inches(1)
                    height = Inches(5)
                    slide.shapes.add_picture(image_stream, left, top, height=height)
                except Exception as e:
                    logger.warning(f"Failed to add image: {e}")
        
        elif slide_type == "chart":
            # 图表页（简单表格）
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加标题
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
            
            # 创建简单的数据表格
            if isinstance(content, list) and len(content) > 0:
                rows = len(content) + 1
                cols = len(content[0]) if isinstance(content[0], list) else 2
                
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                height = Inches(4)
                
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # 设置表头
                header_row = table.rows[0]
                for i in range(cols):
                    header_row.cells[i].text = f"列{i+1}"
                    header_row.cells[i].fill.solid()
                    header_row.cells[i].fill.fore_color.rgb = RGBColor(0, 51, 102)
                    for paragraph in header_row.cells[i].text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                
                # 填充数据
                for i, row_data in enumerate(content):
                    row = table.rows[i+1]
                    if isinstance(row_data, list):
                        for j, cell_data in enumerate(row_data[:cols]):
                            row.cells[j].text = str(cell_data)
                    else:
                        row.cells[0].text = str(row_data)
        
        elif slide_type == "blank":
            # 空白页
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            if title:
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = title
    
    def generate_from_outline(self, title: str, outline: List[Dict[str, Any]]) -> str:
        """
        从大纲生成PPT
        
        Args:
            title: PPT标题
            outline: 大纲结构，每个元素包含:
                - title: 章节标题
                - content: 章节内容（文本或子列表）
        
        Returns:
            生成的PPT文件路径
        """
        slides = []
        
        for section in outline:
            section_title = section.get("title", "")
            section_content = section.get("content", "")
            
            # 添加章节标题页
            slides.append({
                "type": "title",
                "title": section_title
            })
            
            # 添加内容页
            if section_content:
                if isinstance(section_content, list):
                    slides.append({
                        "type": "bullet",
                        "title": section_title,
                        "content": section_content
                    })
                else:
                    slides.append({
                        "type": "content",
                        "title": section_title,
                        "content": section_content
                    })
        
        return self.generate_ppt(title, slides)


class ToolExecutor:
    def __init__(self):
        self.ppt_generator = PPTGenerator()
        self.tools = {
            "execute_script": self._execute_script,
            "read_file": self._read_file,
            "write_file": self._write_file,
            "list_files": self._list_files,
            "web_search": self._web_search,
            "summarize": self._summarize,
            "generate_outline": self._generate_outline,
            "generate_ppt": self._generate_ppt,
            "generate_ppt_from_outline": self._generate_ppt_from_outline,
            "generate_and_send_ppt": self._generate_and_send_ppt,
            "generate_ppt_from_content": self._generate_ppt_from_content
        }
    
    def execute(self, tool_id: str, parameters: Dict[str, Any]) -> str:
        if tool_id not in self.tools:
            return f"Unknown tool: {tool_id}"
        
        try:
            return self.tools[tool_id](parameters)
        except Exception as e:
            logger.error(f"Tool execution failed: {str(e)}")
            return f"Error executing {tool_id}: {str(e)}"
    
    def _execute_script(self, params: Dict[str, Any]) -> str:
        script = params.get("script", "")
        script_type = params.get("type", "python")
        
        output, error, success = sandbox.execute_script(script, script_type)
        
        if success:
            return output
        return f"Script failed: {error}"
    
    def _read_file(self, params: Dict[str, Any]) -> str:
        file_path = params.get("path", "")
        content = sandbox.read_file(file_path)
        return content if content else "Failed to read file"
    
    def _write_file(self, params: Dict[str, Any]) -> str:
        file_path = params.get("path", "")
        content = params.get("content", "")
        
        success = sandbox.write_file(file_path, content)
        return "File written successfully" if success else "Failed to write file"
    
    def _list_files(self, params: Dict[str, Any]) -> str:
        directory = params.get("directory", "")
        files = sandbox.list_files(directory)
        return "\n".join(files) if files else "Failed to list files"
    
    def _web_search(self, params: Dict[str, Any]) -> str:
        query = params.get("query", "")
        return f"Search results for: {query}\n(Web search integration placeholder)"
    
    def _summarize(self, params: Dict[str, Any]) -> str:
        text = params.get("text", "")
        max_length = params.get("max_length", 100)
        
        if len(text) <= max_length:
            return text
        
        return text[:max_length] + "..."
    
    def _generate_ppt(self, params: Dict[str, Any]) -> str:
        """
        生成PPT文件
        
        Parameters:
            title: PPT标题
            slides: 幻灯片列表，每个元素包含:
                - type: 幻灯片类型 (title, content, bullet, image, chart, blank)
                - title: 幻灯片标题
                - content: 内容（文本或列表）
        
        Returns:
            生成的PPT文件路径
        """
        title = params.get("title", "Untitled Presentation")
        slides = params.get("slides", [])
        
        try:
            output_path = self.ppt_generator.generate_ppt(title, slides)
            return f"PPT generated successfully! Path: {output_path}"
        except Exception as e:
            logger.error(f"PPT generation failed: {str(e)}")
            return f"PPT generation failed: {str(e)}"
    
    def _generate_ppt_from_outline(self, params: Dict[str, Any]) -> str:
        """
        从大纲生成PPT
        
        Parameters:
            title: PPT标题
            outline: 大纲结构，每个元素包含:
                - title: 章节标题
                - content: 章节内容（文本或子列表）
        
        Returns:
            生成的PPT文件路径
        """
        title = params.get("title", "Untitled Presentation")
        outline = params.get("outline", [])
        
        try:
            output_path = self.ppt_generator.generate_from_outline(title, outline)
            return f"PPT generated successfully from outline! Path: {output_path}"
        except Exception as e:
            logger.error(f"PPT generation from outline failed: {str(e)}")
            return f"PPT generation failed: {str(e)}"
    
    def _generate_and_send_ppt(self, params: Dict[str, Any]) -> str:
        """
        生成PPT并发送给用户
        
        Parameters:
            user_id: 用户ID
            title: PPT标题
            slides: 幻灯片列表
            im_type: IM类型（feishu, dingtalk, wecom）
        
        Returns:
            操作结果
        """
        import asyncio
        
        user_id = params.get("user_id", "")
        title = params.get("title", "Untitled Presentation")
        slides = params.get("slides", [])
        im_type = params.get("im_type", "feishu")
        
        if not user_id:
            return "Error: user_id is required"
        
        try:
            # 生成PPT
            output_path = self.ppt_generator.generate_ppt(title, slides)
            
            # 获取IM适配器并发送
            from src.plugins import get_im_adapter
            
            adapter = get_im_adapter(im_type)
            if adapter:
                # 发送文件
                import os
                file_name = os.path.basename(output_path)
                send_success = asyncio.run(adapter.send_file(user_id, output_path, file_name))
                
                if send_success:
                    return f"PPT generated and sent successfully! File: {file_name}"
                else:
                    return f"PPT generated but failed to send. Path: {output_path}"
            else:
                return f"PPT generated but IM adapter not available. Path: {output_path}"
                
        except Exception as e:
            logger.error(f"Generate and send PPT failed: {str(e)}")
            return f"Failed to generate and send PPT: {str(e)}"
    
    def _generate_outline(self, params: Dict[str, Any]) -> str:
        """
        根据内容生成PPT大纲
        
        Parameters:
            content: 文件内容或文本
            title: PPT标题（可选）
        
        Returns:
            生成的大纲JSON字符串
        """
        content = params.get("content", "")
        title = params.get("title", "文档总结")
        
        if not content:
            return "Error: content is required"
        
        try:
            from src.infrastructure.model_router import select_model, call_model
            
            model = select_model("summarization", "complex")
            if not model:
                model = select_model("document_analysis", "complex")
            
            if not model:
                return "Error: No suitable model available for outline generation"
            
            prompt = f"""根据以下文档内容，生成一个结构化的PPT大纲：

文档内容：
{content}

要求：
1. 分析文档核心内容，提取关键章节
2. 每个章节包含：章节标题 + 主要要点（3-5个）
3. 大纲结构清晰，逻辑连贯
4. 使用中文输出

PPT标题：{title}

请输出JSON格式，结构如下：
[
  {{"title": "章节1标题", "content": ["要点1", "要点2", "要点3"]}},
  {{"title": "章节2标题", "content": ["要点1", "要点2"]}}
]
"""
            result = call_model(model, [{"role": "user", "content": prompt}])
            
            import json
            try:
                outline = json.loads(result)
                if isinstance(outline, list) and len(outline) > 0:
                    logger.info(f"大纲生成成功，共 {len(outline)} 个章节")
                    return json.dumps(outline, ensure_ascii=False, indent=2)
                else:
                    return f"Error: Invalid outline format: {result}"
            except json.JSONDecodeError:
                # 如果不是JSON格式，尝试解析为列表
                return f"Error: Failed to parse outline: {result[:200]}..."
                
        except Exception as e:
            logger.error(f"Outline generation failed: {str(e)}")
            return f"大纲生成失败: {str(e)}"
    
    def _get_content_from_vector_db(self, file_key: str) -> Optional[str]:
        """
        从向量数据库中获取文件内容
        
        Args:
            file_key: 文件标识
        
        Returns:
            文件内容字符串，如果未找到返回None
        """
        try:
            from src.data.database import db
            
            # 从记忆存储中查询包含该file_key的记录
            memories = db.get_memories_by_tag(file_key)
            
            if memories:
                # 返回最新的内容
                memories.sort(key=lambda m: m.timestamp, reverse=True)
                return memories[0].content
            
            logger.info(f"未在向量数据库中找到文件: {file_key}")
            return None
            
        except Exception as e:
            logger.error(f"从向量数据库获取文件内容失败: {str(e)}")
            return None
    
    def _generate_ppt_from_content(self, params: Dict[str, Any]) -> str:
        """
        根据文件内容直接生成PPT（完整流程：内容→大纲→PPT）
        
        Parameters:
            content: 文件内容（直接提供内容）
            file_key: 飞书文件key（通过file_key从本地存储获取内容）
            title: PPT标题（可选）
            user_id: 用户ID（用于用户隔离）
        
        Returns:
            生成的PPT文件路径
        """
        content = params.get("content", "")
        title = params.get("title", "文档总结")
        file_key = params.get("file_key", "")
        user_id = params.get("user_id", "")
        
        # 如果没有直接提供content，但提供了file_key，则先从本地存储获取文件内容
        if not content and file_key:
            logger.info(f"通过file_key从本地存储获取文件内容: {file_key}")
            
            # 首先尝试从向量数据库中获取（文件上传时已存储，支持用户隔离）
            content = self._get_content_from_vector_db(file_key, user_id)
            
            # 如果本地没有，再尝试从飞书下载
            if not content:
                logger.info(f"本地存储未找到，尝试从飞书下载: {file_key}")
                from src.plugins.im_adapters import FeishuAdapter
                feishu_adapter = FeishuAdapter()
                try:
                    content = feishu_adapter.read_file(file_key)
                    if not content:
                        return "Error: Failed to read file content"
                    logger.info(f"从飞书下载文件成功，内容长度: {len(content)}")
                except Exception as e:
                    logger.error(f"从飞书读取文件失败: {str(e)}")
                    return f"Error: Failed to read file: {str(e)}"
            else:
                logger.info(f"从本地存储获取文件内容成功，内容长度: {len(content)}")
        
        if not content:
            return "Error: content is required"
        
        try:
            logger.info(f"开始根据内容生成PPT，内容长度: {len(content)}")
            
            # 步骤1：根据内容生成大纲
            outline_result = self._generate_outline({"content": content, "title": title})
            
            import json
            try:
                outline = json.loads(outline_result)
                if not isinstance(outline, list) or len(outline) == 0:
                    return f"Error: Failed to generate outline: {outline_result}"
            except json.JSONDecodeError:
                return f"Error: Invalid outline: {outline_result[:100]}..."
            
            logger.info(f"大纲生成成功，共 {len(outline)} 个章节")
            
            # 步骤2：根据大纲生成PPT
            output_path = self.ppt_generator.generate_from_outline(title, outline)
            
            logger.info(f"PPT生成成功: {output_path}")
            return f"PPT生成成功！\n\n大纲章节: {len(outline)} 个\n文件路径: {output_path}"
            
        except Exception as e:
            logger.error(f"Generate PPT from content failed: {str(e)}")
            return f"根据内容生成PPT失败: {str(e)}"


tool_executor = ToolExecutor()
