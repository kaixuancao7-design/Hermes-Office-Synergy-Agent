"""质量门控 - PPT质量自动检查"""

import os
import re
from typing import Dict, Any, List, Optional, Callable
from dataclasses import dataclass, field
from src.logging_config import get_logger

logger = get_logger("engine.quality_gate")


@dataclass
class QualityError:
    """质量错误"""
    level: str  # "error" or "warning"
    check_name: str
    message: str
    location: Optional[str] = None


@dataclass
class QualityResult:
    """质量检查结果"""
    passed: bool
    errors: List[QualityError] = field(default_factory=list)
    warnings: List[QualityError] = field(default_factory=list)

    def add_error(self, check_name: str, message: str, location: str = None):
        self.errors.append(QualityError("error", check_name, message, location))
        self.passed = False

    def add_warning(self, check_name: str, message: str, location: str = None):
        self.warnings.append(QualityError("warning", check_name, message, location))


class QualityGate:
    """质量门控 - PPT质量自动检查

    在PPT生成过程中和生成后进行质量检查，确保输出符合设计规范。
    错误必须修复才能继续，警告建议修复但可以跳过。
    """

    def __init__(self, strict_mode: bool = True):
        """
        初始化质量门控

        Args:
            strict_mode: 严格模式，True时警告也会阻止流程
        """
        self.strict_mode = strict_mode
        self._checks: Dict[str, Callable] = {}
        self._register_default_checks()

    def _register_default_checks(self):
        """注册默认检查项"""
        self._checks = {
            "file_exists": self._check_file_exists,
            "file_size": self._check_file_size,
            "pptx_structure": self._check_pptx_structure,
            "slide_count": self._check_slide_count,
            "title_present": self._check_title_present,
            "font_safety": self._check_font_safety,
        }

    def gate(self, ppt_path: str, **kwargs) -> QualityResult:
        """
        执行质量门控检查

        Args:
            ppt_path: PPT文件路径
            **kwargs: 额外参数（如min_slides, max_slides等）

        Returns:
            质量检查结果
        """
        result = QualityResult(passed=True)

        logger.info(f"QualityGate started for: {ppt_path}")

        for check_name, check_func in self._checks.items():
            try:
                check_result = check_func(ppt_path, **kwargs)
                if not check_result["passed"]:
                    if check_result.get("level") == "error":
                        result.add_error(
                            check_name,
                            check_result.get("message", "Check failed"),
                            check_result.get("location")
                        )
                    else:
                        result.add_warning(
                            check_name,
                            check_result.get("message", "Check warning"),
                            check_result.get("location")
                        )
            except Exception as e:
                logger.error(f"Quality check {check_name} failed with exception: {str(e)}")
                result.add_error(check_name, f"检查异常: {str(e)}")

        if result.passed:
            logger.info(f"QualityGate passed: {ppt_path}")
        else:
            logger.warning(f"QualityGate failed: {ppt_path}, errors={len(result.errors)}, warnings={len(result.warnings)}")

        return result

    def _check_file_exists(self, ppt_path: str, **kwargs) -> Dict[str, Any]:
        """检查文件是否存在"""
        if not os.path.exists(ppt_path):
            return {
                "passed": False,
                "level": "error",
                "message": f"PPT文件不存在: {ppt_path}"
            }
        return {"passed": True}

    def _check_file_size(self, ppt_path: str, **kwargs) -> Dict[str, Any]:
        """检查文件大小是否合理"""
        min_size = kwargs.get("min_size", 1024)  # 最少1KB
        max_size = kwargs.get("max_size", 50 * 1024 * 1024)  # 最多50MB

        try:
            size = os.path.getsize(ppt_path)
            if size < min_size:
                return {
                    "passed": False,
                    "level": "error",
                    "message": f"PPT文件过小，可能生成失败: {size} bytes"
                }
            if size > max_size:
                return {
                    "passed": False,
                    "level": "warning",
                    "message": f"PPT文件过大: {size / (1024*1024):.2f} MB"
                }
        except Exception as e:
            return {
                "passed": False,
                "level": "error",
                "message": f"无法获取文件大小: {str(e)}"
            }

        return {"passed": True}

    def _check_pptx_structure(self, ppt_path: str, **kwargs) -> Dict[str, Any]:
        """检查PPTX文件结构是否完整"""
        try:
            from pptx import Presentation
            from zipfile import ZipFile

            if not zipfile.is_zipfile(ppt_path):
                return {
                    "passed": False,
                    "level": "error",
                    "message": "PPTX文件格式无效（不是有效的ZIP文件）"
                }

            prs = Presentation(ppt_path)

            return {"passed": True}

        except ImportError:
            logger.warning("python-pptx not available, skipping structure check")
            return {"passed": True}
        except Exception as e:
            return {
                "passed": False,
                "level": "error",
                "message": f"PPTX文件结构检查失败: {str(e)}"
            }

    def _check_slide_count(self, ppt_path: str, **kwargs) -> Dict[str, Any]:
        """检查幻灯片数量是否在合理范围内"""
        min_slides = kwargs.get("min_slides", 1)
        max_slides = kwargs.get("max_slides", 100)

        try:
            from pptx import Presentation
            prs = Presentation(ppt_path)
            slide_count = len(prs.slides)

            if slide_count < min_slides:
                return {
                    "passed": False,
                    "level": "error",
                    "message": f"幻灯片数量过少: {slide_count} < {min_slides}"
                }
            if slide_count > max_slides:
                return {
                    "passed": False,
                    "level": "warning",
                    "message": f"幻灯片数量过多: {slide_count} > {max_slides}"
                }

            return {"passed": True}

        except ImportError:
            logger.warning("python-pptx not available, skipping slide count check")
            return {"passed": True}
        except Exception as e:
            return {
                "passed": False,
                "level": "error",
                "message": f"幻灯片数量检查失败: {str(e)}"
            }

    def _check_title_present(self, ppt_path: str, **kwargs) -> Dict[str, Any]:
        """检查是否包含标题"""
        try:
            from pptx import Presentation
            prs = Presentation(ppt_path)

            has_title = False
            for slide in prs.slides:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    has_title = True
                    break

            if not has_title:
                return {
                    "passed": False,
                    "level": "warning",
                    "message": "PPT中没有发现标题"
                }

            return {"passed": True}

        except ImportError:
            return {"passed": True}
        except Exception as e:
            return {
                "passed": False,
                "level": "warning",
                "message": f"标题检查失败: {str(e)}"
            }

    def _check_font_safety(self, ppt_path: str, **kwargs) -> Dict[str, Any]:
        """检查字体是否PPT安全（常见字体）"""
        safe_fonts = [
            "Arial", "Calibri", "Times New Roman", "Microsoft YaHei",
            "SimSun", "SimHei", "FangSong", "KaiTi"
        ]

        try:
            from pptx import Presentation
            from pptx.util import Pt

            prs = Presentation(ppt_path)
            unsafe_fonts = set()

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name and run.font.name not in safe_fonts:
                                    unsafe_fonts.add(run.font.name)

            if unsafe_fonts:
                return {
                    "passed": True,
                    "level": "warning",
                    "message": f"使用了非标准字体: {', '.join(unsafe_fonts)}"
                }

            return {"passed": True}

        except ImportError:
            return {"passed": True}
        except Exception as e:
            return {
                "passed": True,
                "level": "warning",
                "message": f"字体检查失败: {str(e)}"
            }

    def add_check(self, name: str, check_func: Callable):
        """添加自定义检查项"""
        self._checks[name] = check_func
        logger.info(f"Custom check added: {name}")

    def remove_check(self, name: str):
        """移除检查项"""
        if name in self._checks:
            del self._checks[name]
            logger.info(f"Check removed: {name}")

    def is_blocking(self, result: QualityResult) -> bool:
        """判断结果是否应该阻止流程"""
        if not result.passed:
            return True
        if self.strict_mode and result.warnings:
            return True
        return False

    def format_report(self, result: QualityResult) -> str:
        """格式化质量报告"""
        lines = []

        if result.passed and not result.warnings:
            lines.append("✅ 质量检查通过")
        elif result.passed and result.warnings:
            lines.append("⚠️ 质量检查通过，但有警告:")
        else:
            lines.append("❌ 质量检查失败:")

        for error in result.errors:
            lines.append(f"  🔴 [{error.check_name}] {error.message}")

        for warning in result.warnings:
            lines.append(f"  🟡 [{warning.check_name}] {warning.message}")

        return "\n".join(lines)


import zipfile

quality_gate = QualityGate()
