"""统一文档加载器 - 提供统一的文档加载和预处理管道"""
import os
import re
from typing import List, Dict, Any, Optional, Tuple
from abc import ABC, abstractmethod

from src.config import settings
from src.logging_config import get_logger

logger = get_logger("document_loader")


class DocumentLoaderBase(ABC):
    """文档加载器基类"""
    
    @abstractmethod
    def load(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """加载文档"""
        pass
    
    @abstractmethod
    def supports(self, file_extension: str) -> bool:
        """检查是否支持该文件类型"""
        pass


class TxtLoader(DocumentLoaderBase):
    """纯文本文件加载器"""
    
    def load(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            encoding = kwargs.get("encoding", "utf-8")
            with open(file_path, "r", encoding=encoding) as f:
                content = f.read()
            return {
                "success": True,
                "content": content,
                "metadata": {"file_type": "txt", "encoding": encoding}
            }
        except Exception as e:
            logger.error(f"加载TXT文件失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [".txt"]


class MarkdownLoader(DocumentLoaderBase):
    """Markdown文件加载器"""
    
    def load(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            encoding = kwargs.get("encoding", "utf-8")
            with open(file_path, "r", encoding=encoding) as f:
                content = f.read()
            return {
                "success": True,
                "content": content,
                "metadata": {"file_type": "markdown", "encoding": encoding}
            }
        except Exception as e:
            logger.error(f"加载Markdown文件失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [".md", ".markdown"]


class DocxLoader(DocumentLoaderBase):
    """DOCX文件加载器"""
    
    def load(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            from docx import Document
            doc = Document(file_path)
            content = "\n".join([para.text for para in doc.paragraphs])
            return {
                "success": True,
                "content": content,
                "metadata": {"file_type": "docx", "paragraph_count": len(doc.paragraphs)}
            }
        except ImportError:
            logger.warning("python-docx库未安装")
            return {"success": False, "error": "python-docx库未安装", "content": "", "metadata": {}}
        except Exception as e:
            logger.error(f"加载DOCX文件失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [".docx"]


class XlsxLoader(DocumentLoaderBase):
    """XLSX文件加载器"""
    
    def load(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            import pandas as pd
            df = pd.read_excel(file_path, sheet_name=None)
            
            content = ""
            for sheet_name, sheet_data in df.items():
                content += f"=== 工作表: {sheet_name} ===\n"
                content += sheet_data.to_string(index=False)
                content += "\n\n"
            
            return {
                "success": True,
                "content": content.strip(),
                "metadata": {
                    "file_type": "xlsx",
                    "sheets": list(df.keys()),
                    "total_sheets": len(df)
                }
            }
        except ImportError:
            logger.warning("pandas库未安装")
            return {"success": False, "error": "pandas库未安装", "content": "", "metadata": {}}
        except Exception as e:
            logger.error(f"加载XLSX文件失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [".xlsx", ".xls"]


class PdfLoader(DocumentLoaderBase):
    """PDF文件加载器"""
    
    def load(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            content = "\n".join([page.extract_text() for page in reader.pages])
            return {
                "success": True,
                "content": content,
                "metadata": {"file_type": "pdf", "page_count": len(reader.pages)}
            }
        except ImportError:
            logger.warning("PyPDF2库未安装")
            return {"success": False, "error": "PyPDF2库未安装", "content": "", "metadata": {}}
        except Exception as e:
            logger.error(f"加载PDF文件失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [".pdf"]


class PptxLoader(DocumentLoaderBase):
    """PPTX文件加载器"""
    
    def load(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            content = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                content += f"=== 幻灯片 {slide_num} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
                content += "\n"
            
            return {
                "success": True,
                "content": content.strip(),
                "metadata": {"file_type": "pptx", "slide_count": len(prs.slides)}
            }
        except ImportError:
            logger.warning("python-pptx库未安装")
            return {"success": False, "error": "python-pptx库未安装", "content": "", "metadata": {}}
        except Exception as e:
            logger.error(f"加载PPTX文件失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [".pptx"]


class HtmlLoader(DocumentLoaderBase):
    """HTML文件加载器"""
    
    def load(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f, "html.parser")
            
            content = soup.get_text(separator="\n")
            return {
                "success": True,
                "content": content,
                "metadata": {"file_type": "html"}
            }
        except ImportError:
            logger.warning("beautifulsoup4库未安装")
            return {"success": False, "error": "beautifulsoup4库未安装", "content": "", "metadata": {}}
        except Exception as e:
            logger.error(f"加载HTML文件失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [".html", ".htm"]


class JsonLoader(DocumentLoaderBase):
    """JSON文件加载器"""
    
    def load(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            import json
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            
            content = json.dumps(data, ensure_ascii=False, indent=2)
            return {
                "success": True,
                "content": content,
                "metadata": {"file_type": "json"}
            }
        except Exception as e:
            logger.error(f"加载JSON文件失败: {str(e)}")
            return {"success": False, "error": str(e), "content": "", "metadata": {}}
    
    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [".json"]


class DocumentLoader:
    """统一文档加载器"""
    
    def __init__(self):
        self.loaders: List[DocumentLoaderBase] = [
            TxtLoader(),
            MarkdownLoader(),
            DocxLoader(),
            XlsxLoader(),
            PdfLoader(),
            PptxLoader(),
            HtmlLoader(),
            JsonLoader()
        ]
    
    def load(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """加载文档"""
        _, ext = os.path.splitext(file_path)
        
        for loader in self.loaders:
            if loader.supports(ext):
                return loader.load(file_path, **kwargs)
        
        return {
            "success": False,
            "error": f"不支持的文件类型: {ext}",
            "content": "",
            "metadata": {}
        }
    
    def load_from_content(self, content: bytes, file_name: str, **kwargs) -> Dict[str, Any]:
        """从内容加载文档"""
        import tempfile
        
        _, ext = os.path.splitext(file_name)
        
        for loader in self.loaders:
            if loader.supports(ext):
                with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as f:
                    f.write(content)
                    temp_path = f.name
                
                try:
                    result = loader.load(temp_path, **kwargs)
                    result["metadata"]["original_filename"] = file_name
                    return result
                finally:
                    os.unlink(temp_path)
        
        return {
            "success": False,
            "error": f"不支持的文件类型: {ext}",
            "content": "",
            "metadata": {"original_filename": file_name}
        }
    
    def get_supported_types(self) -> List[str]:
        """获取支持的文件类型"""
        types = set()
        for loader in self.loaders:
            # 检查每种加载器支持的类型
            for ext in [".txt", ".md", ".docx", ".xlsx", ".xls", ".pdf", ".pptx", ".html", ".htm", ".json"]:
                if loader.supports(ext):
                    types.add(ext)
        return sorted(list(types))


class DocumentPreprocessor:
    """文档预处理器"""
    
    def __init__(self):
        self.clean_rules = [
            self._remove_extra_whitespace,
            self._remove_special_chars,
            self._normalize_unicode,
            self._remove_empty_lines
        ]
    
    def process(self, content: str, operations: Optional[List[str]] = None) -> str:
        """处理文档内容"""
        if not content:
            return ""
        
        # 确保content是字符串类型
        if not isinstance(content, str):
            content = str(content)
        
        # 默认执行所有清洗操作
        if operations is None:
            operations = ["whitespace", "special_chars", "unicode", "empty_lines"]
        
        if "whitespace" in operations:
            content = self._remove_extra_whitespace(content)
        
        if "special_chars" in operations:
            content = self._remove_special_chars(content)
        
        if "unicode" in operations:
            content = self._normalize_unicode(content)
        
        if "empty_lines" in operations:
            content = self._remove_empty_lines(content)
        
        if "lowercase" in operations:
            content = content.lower()
        
        if "trim" in operations:
            content = content.strip()
        
        return content
    
    def _remove_extra_whitespace(self, content: str) -> str:
        """移除多余空白字符"""
        content = re.sub(r'\s+', ' ', content)
        return content.strip()
    
    def _remove_special_chars(self, content: str) -> str:
        """移除特殊字符"""
        # 保留基本标点符号
        content = re.sub(r'[^\w\s\u4e00-\u9fff，。！？；：、""''（）【】《》<>{}]', '', content)
        return content
    
    def _normalize_unicode(self, content: str) -> str:
        """标准化Unicode字符"""
        try:
            if hasattr(content, 'normalize'):
                return content.normalize('NFKC')
            else:
                # 如果字符串没有normalize方法，跳过此步骤
                return content
        except Exception:
            return content
    
    def _remove_empty_lines(self, content: str) -> str:
        """移除空行"""
        lines = content.split('\n')
        lines = [line.strip() for line in lines if line.strip()]
        return '\n'.join(lines)
    
    def chunk_content(self, content: str, chunk_size: int = 500, 
                     chunk_overlap: int = 50) -> List[str]:
        """将内容分割为块"""
        chunks = []
        start = 0
        
        while start < len(content):
            end = min(start + chunk_size, len(content))
            
            # 在合适的位置分割
            if end < len(content):
                # 优先在段落边界分割
                split_pos = content.rfind('\n', start, end)
                if split_pos == -1 or split_pos < start + chunk_size * 0.8:
                    split_pos = content.rfind('。', start, end)
                if split_pos == -1 or split_pos < start + chunk_size * 0.8:
                    split_pos = content.rfind('！', start, end)
                if split_pos == -1 or split_pos < start + chunk_size * 0.8:
                    split_pos = content.rfind('？', start, end)
                if split_pos == -1:
                    split_pos = end - 1
                
                end = split_pos + 1
            
            chunk = content[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # 计算下一个起始位置（考虑重叠）
            next_start = end - chunk_overlap
            if next_start <= start:
                next_start = end
            
            start = next_start
        
        return chunks


class DocumentPipeline:
    """文档处理管道"""
    
    def __init__(self):
        self.loader = DocumentLoader()
        self.preprocessor = DocumentPreprocessor()
    
    def process(self, file_path: str = None, content: bytes = None, 
                file_name: str = None, **kwargs) -> Dict[str, Any]:
        """
        完整的文档处理流程
        
        Args:
            file_path: 文件路径
            content: 文件内容（字节）
            file_name: 文件名（当使用content参数时必需）
            kwargs: 额外参数
        
        Returns:
            处理结果
        """
        # 1. 加载文档
        if file_path:
            load_result = self.loader.load(file_path, **kwargs)
        elif content and file_name:
            load_result = self.loader.load_from_content(content, file_name, **kwargs)
        else:
            return {"success": False, "error": "必须提供file_path或content+file_name"}
        
        if not load_result["success"]:
            return load_result
        
        # 2. 预处理
        operations = kwargs.get("operations", ["whitespace", "special_chars", "unicode", "empty_lines"])
        processed_content = self.preprocessor.process(load_result["content"], operations)
        
        # 3. 分割（可选）
        chunk_size = kwargs.get("chunk_size", 0)
        if chunk_size > 0:
            chunks = self.preprocessor.chunk_content(
                processed_content, 
                chunk_size=chunk_size, 
                chunk_overlap=kwargs.get("chunk_overlap", 50)
            )
            load_result["chunks"] = chunks
            load_result["chunk_count"] = len(chunks)
        
        load_result["processed_content"] = processed_content
        
        return load_result


# 全局实例
document_loader = DocumentLoader()
document_preprocessor = DocumentPreprocessor()
document_pipeline = DocumentPipeline()


# 便捷函数
def load_document(file_path: str, **kwargs) -> Dict[str, Any]:
    """加载文档"""
    return document_loader.load(file_path, **kwargs)


def process_document(file_path: str = None, content: bytes = None, 
                    file_name: str = None, **kwargs) -> Dict[str, Any]:
    """处理文档（加载+预处理）"""
    return document_pipeline.process(file_path, content, file_name, **kwargs)
