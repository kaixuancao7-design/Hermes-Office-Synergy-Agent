"""工具执行插件实现"""
import os
import re
from typing import Dict, Any, List, Optional
from urllib.parse import unquote

from src.plugins.base import ToolExecutorBase
from src.config import settings
from src.utils import generate_id, get_timestamp
from src.logging_config import get_logger

logger = get_logger("tool")


def decode_filename(content_disposition: str) -> str:
    """
    从 Content-Disposition 响应头中提取并解码文件名
    支持 RFC 5987 编码（filename*=UTF-8''xxx）和普通格式
    处理中文文件名编码问题
    
    Args:
        content_disposition: HTTP响应头中的 Content-Disposition 值
    
    Returns:
        解码后的文件名，如果无法提取则返回空字符串
    """
    if not content_disposition:
        return ""
    
    # 优先尝试解析 filename*=UTF-8''xxx 格式（RFC 5987 编码）
    utf8_match = re.search(r"filename\*=UTF-8''([^\s;]+)", content_disposition, re.IGNORECASE)
    if utf8_match:
        try:
            return unquote(utf8_match.group(1), encoding='utf-8')
        except:
            pass
    
    # 尝试解析普通 filename=xxx 格式
    match = re.search(r'filename[^;=\n]*=(([""]).*?\2|[^;\n]*)', content_disposition)
    if match:
        file_name = match.group(1).strip('"')
        # 处理可能的编码问题：ISO-8859-1 编码的中文
        try:
            # 先尝试直接使用（检查是否是有效的UTF-8）
            file_name.encode('utf-8')
        except UnicodeEncodeError:
            # 如果不是有效的UTF-8，尝试ISO-8859-1解码
            try:
                file_name = file_name.encode('ISO-8859-1').decode('utf-8')
            except:
                # 解码失败，保持原样
                pass
        return file_name
    
    return ""


class BasicToolExecutor(ToolExecutorBase):
    """基础工具执行器"""
    
    def __init__(self):
        self.tools = {}
        self.tool_name_mapping = {
            "file_reader": "feishu_file_read",
            "read_file": "feishu_file_read",
            "feishu_reader": "feishu_file_read"
        }
        
        # 初始化飞书客户端
        self._initialize_feishu_client()
        
        # 初始化记忆存储
        self._initialize_memory_store()
        
        self._register_default_tools()
    
    def _initialize_feishu_client(self):
        """初始化飞书API客户端"""
        feishu_app_id = os.getenv("FEISHU_APP_ID")
        feishu_app_secret = os.getenv("FEISHU_APP_SECRET")
        
        if feishu_app_id and feishu_app_secret:
            try:
                from lark_oapi import Client, FEISHU_DOMAIN
                self._feishu_client = Client.builder() \
                    .app_id(feishu_app_id) \
                    .app_secret(feishu_app_secret) \
                    .domain(FEISHU_DOMAIN) \
                    .build()
                logger.info("飞书API客户端初始化成功")
            except Exception as e:
                self._feishu_client = None
                logger.error(f"飞书API客户端初始化失败: {str(e)}")
        else:
            self._feishu_client = None
            logger.warning("飞书API客户端未初始化：缺少 FEISHU_APP_ID 或 FEISHU_APP_SECRET")
    
    def _initialize_memory_store(self):
        """初始化记忆存储"""
        try:
            from src.engine.memory_manager import MemoryManager
            self.memory_manager = MemoryManager()
            logger.info("记忆存储初始化成功")
        except Exception as e:
            self.memory_manager = None
            logger.warning(f"记忆存储未初始化：{str(e)}")
    
    def _register_default_tools(self):
        """注册默认工具"""
        self.register_tool("document_search", DocumentSearchTool)
        self.register_tool("memory_search", MemorySearchTool)
        self.register_tool("web_search", WebSearchTool)
        self.register_tool("code_execution", CodeExecutionTool)
        self.register_tool("file_operations", FileOperationsTool)
        self.register_tool("feishu_file_read", FeishuFileReadTool)
        # PPT生成工具
        self.register_tool("generate_outline", GenerateOutlineTool)
        self.register_tool("generate_ppt", GeneratePPTFromSlidesTool)
        self.register_tool("generate_ppt_from_outline", GeneratePPTFromOutlineTool)
        self.register_tool("generate_ppt_from_content", GeneratePPTFromContentTool)
    
    def execute(self, tool_id: str, parameters: Dict[str, Any]) -> Dict[str, Any]:
        """执行工具"""
        # 工具名称映射
        if tool_id in self.tool_name_mapping:
            tool_id = self.tool_name_mapping[tool_id]
        
        if tool_id not in self.tools:
            return {"success": False, "error": f"工具 {tool_id} 不存在"}
        
        try:
            # 直接使用传入的参数（ReAct引擎已经处理好了参数格式）
            # 注意：不再尝试提取嵌套的parameters，因为ReAct引擎已经处理过了
            tool_params = parameters or {}
            
            # 获取工具类并创建实例（传入executor，以便工具可以访问客户端和存储）
            tool_class = self.tools[tool_id]
            tool_instance = tool_class(executor=self)
            result = tool_instance.execute(tool_params)
            return result
        except Exception as e:
            logger.error(f"执行工具 {tool_id} 失败: {str(e)}")
            return {"success": False, "error": str(e)}
    
    def register_tool(self, tool_id: str, tool_class: Any) -> bool:
        """注册工具"""
        self.tools[tool_id] = tool_class
        logger.info(f"注册工具: {tool_id}")
        return True
    
    def get_tools(self) -> List[str]:
        """获取所有可用工具"""
        return list(self.tools.keys())
    
    def get_executor_type(self) -> str:
        return "basic"


class SandboxedToolExecutor(ToolExecutorBase):
    """沙箱工具执行器（带安全限制）"""
    
    def __init__(self):
        self.tools = {}
        self.allowed_paths = settings.SANDBOX_ALLOWED_PATHS or []
        self.max_execution_time = settings.SANDBOX_MAX_EXECUTION_TIME or 30
        
        # 初始化飞书客户端
        self._initialize_feishu_client()
        
        # 初始化记忆存储
        self._initialize_memory_store()
        
        self._register_default_tools()
    
    def _initialize_feishu_client(self):
        """初始化飞书API客户端"""
        feishu_app_id = os.getenv("FEISHU_APP_ID")
        feishu_app_secret = os.getenv("FEISHU_APP_SECRET")
        
        if feishu_app_id and feishu_app_secret:
            try:
                from lark_oapi import Client, FEISHU_DOMAIN
                self._feishu_client = Client.builder() \
                    .app_id(feishu_app_id) \
                    .app_secret(feishu_app_secret) \
                    .domain(FEISHU_DOMAIN) \
                    .build()
                logger.info("飞书API客户端初始化成功")
            except Exception as e:
                self._feishu_client = None
                logger.error(f"飞书API客户端初始化失败: {str(e)}")
        else:
            self._feishu_client = None
            logger.warning("飞书API客户端未初始化：缺少 FEISHU_APP_ID 或 FEISHU_APP_SECRET")
    
    def _initialize_memory_store(self):
        """初始化记忆存储"""
        try:
            from src.engine.memory_manager import MemoryManager
            self.memory_manager = MemoryManager()
            logger.info("记忆存储初始化成功")
        except Exception as e:
            self.memory_manager = None
            logger.warning(f"记忆存储未初始化：{str(e)}")
    
    def _register_default_tools(self):
        """注册默认工具"""
        self.register_tool("document_search", DocumentSearchTool)
        self.register_tool("memory_search", MemorySearchTool)
        self.register_tool("web_search", WebSearchTool)
        self.register_tool("code_execution", SandboxedCodeExecutionTool)
        self.register_tool("file_operations", SandboxedFileOperationsTool)
        self.register_tool("feishu_file_read", FeishuFileReadTool)
        # PPT生成工具
        self.register_tool("generate_outline", GenerateOutlineTool)
        self.register_tool("generate_ppt", GeneratePPTFromSlidesTool)
        self.register_tool("generate_ppt_from_outline", GeneratePPTFromOutlineTool)
        self.register_tool("generate_ppt_from_content", GeneratePPTFromContentTool)
    
    def execute(self, tool_id: str, parameters: Dict[str, Any]) -> Dict[str, Any]:
        """执行工具（带安全检查）"""
        if tool_id not in self.tools:
            return {"success": False, "error": f"工具 {tool_id} 不存在"}
        
        # 安全检查
        if not self._check_security(tool_id, parameters):
            return {"success": False, "error": "安全检查失败"}
        
        try:
            # 直接使用传入的参数（ReAct引擎已经处理好了参数格式）
            # 如果参数中包含嵌套的parameters字段，则使用嵌套的参数
            tool_params = parameters.get("parameters", parameters)
            
            # 如果是None，转换为空字典
            if tool_params is None:
                tool_params = {}
            
            # 创建工具实例（传入executor，以便工具可以访问客户端和存储）
            tool_instance = self.tools[tool_id](executor=self)
            result = tool_instance.execute(tool_params)
            return result
        except Exception as e:
            logger.error(f"执行工具 {tool_id} 失败: {str(e)}")
            return {"success": False, "error": str(e)}
    
    def _check_security(self, tool_id: str, parameters: Dict[str, Any]) -> bool:
        """安全检查"""
        if tool_id == "file_operations":
            file_path = parameters.get("path", "")
            if not self._is_path_allowed(file_path):
                logger.warning(f"非法路径访问: {file_path}")
                return False
        
        return True
    
    def _is_path_allowed(self, file_path: str) -> bool:
        """检查路径是否在白名单中"""
        if not file_path:
            return False
        
        for allowed in self.allowed_paths:
            if file_path.startswith(allowed):
                return True
        
        return False
    
    def register_tool(self, tool_id: str, tool_class: Any) -> bool:
        """注册工具"""
        self.tools[tool_id] = tool_class
        logger.info(f"注册工具: {tool_id}")
        return True
    
    def get_tools(self) -> List[str]:
        """获取所有可用工具"""
        return list(self.tools.keys())
    
    def get_executor_type(self) -> str:
        return "sandboxed"


# 工具类定义
class DocumentSearchTool:
    """文档搜索工具"""
    
    def __init__(self, executor=None):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        query = parameters.get("query", "")
        limit = parameters.get("limit", 5)
        
        try:
            # 文档搜索功能暂未实现，返回空结果
            logger.warning("Document search not implemented, returning empty results")
            return {
                "success": True,
                "result": []
            }
        except Exception as e:
            return {"success": False, "error": str(e)}


class MemorySearchTool:
    """记忆搜索工具"""
    
    def __init__(self, executor=None):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        user_id = parameters.get("user_id", "")
        query = parameters.get("query", "")
        limit = parameters.get("limit", 5)
        
        try:
            from src.engine.memory_manager import memory_manager
            results = memory_manager.search_long_term_memory(user_id, query, limit)
            
            return {
                "success": True,
                "result": [r.model_dump() for r in results]
            }
        except Exception as e:
            return {"success": False, "error": str(e)}


class WebSearchTool:
    """网页搜索工具"""
    
    def __init__(self, executor=None):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        query = parameters.get("query", "")
        
        try:
            # 模拟网页搜索结果
            results = [
                {"title": "搜索结果1", "url": "https://example.com", "summary": f"关于 '{query}' 的信息..."},
                {"title": "搜索结果2", "url": "https://example.org", "summary": f"更多关于 '{query}' 的内容..."}
            ]
            
            return {
                "success": True,
                "result": results
            }
        except Exception as e:
            return {"success": False, "error": str(e)}


class CodeExecutionTool:
    """代码执行工具"""
    
    def __init__(self, executor=None):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        code = parameters.get("code", "")
        language = parameters.get("language", "python")
        
        try:
            if language == "python":
                import subprocess
                import tempfile
                import os
                
                with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False) as f:
                    f.write(code)
                    temp_path = f.name
                
                result = subprocess.run(
                    ["python", temp_path],
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                
                os.unlink(temp_path)
                
                return {
                    "success": True,
                    "result": {
                        "stdout": result.stdout,
                        "stderr": result.stderr,
                        "return_code": result.returncode
                    }
                }
            else:
                return {"success": False, "error": f"不支持的语言: {language}"}
        
        except subprocess.TimeoutExpired:
            return {"success": False, "error": "执行超时"}
        except Exception as e:
            return {"success": False, "error": str(e)}


class SandboxedCodeExecutionTool:
    """沙箱代码执行工具"""
    
    def __init__(self, executor: "SandboxedToolExecutor"):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        code = parameters.get("code", "")
        language = parameters.get("language", "python")
        
        # 安全检查：禁止危险操作
        dangerous_patterns = ["import os", "import subprocess", "open(", "os.system", "__import__"]
        for pattern in dangerous_patterns:
            if pattern in code:
                return {"success": False, "error": "禁止执行危险操作"}
        
        try:
            if language == "python":
                # 限制执行时间
                import subprocess
                import tempfile
                import os
                
                with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False) as f:
                    f.write(code)
                    temp_path = f.name
                
                result = subprocess.run(
                    ["python", temp_path],
                    capture_output=True,
                    text=True,
                    timeout=self.executor.max_execution_time
                )
                
                os.unlink(temp_path)
                
                return {
                    "success": True,
                    "result": {
                        "stdout": result.stdout,
                        "stderr": result.stderr,
                        "return_code": result.returncode
                    }
                }
            else:
                return {"success": False, "error": f"不支持的语言: {language}"}
        
        except subprocess.TimeoutExpired:
            return {"success": False, "error": "执行超时"}
        except Exception as e:
            return {"success": False, "error": str(e)}


class FileOperationsTool:
    """文件操作工具"""
    
    def __init__(self, executor=None):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        operation = parameters.get("operation", "")
        path = parameters.get("path", "")
        content = parameters.get("content", "")
        
        try:
            if operation == "read":
                with open(path, "r", encoding="utf-8") as f:
                    return {"success": True, "result": f.read()}
            
            elif operation == "write":
                with open(path, "w", encoding="utf-8") as f:
                    f.write(content)
                return {"success": True, "result": "文件写入成功"}
            
            elif operation == "append":
                with open(path, "a", encoding="utf-8") as f:
                    f.write(content)
                return {"success": True, "result": "内容追加成功"}
            
            elif operation == "list":
                import os
                files = os.listdir(path)
                return {"success": True, "result": files}
            
            else:
                return {"success": False, "error": f"未知操作: {operation}"}
        
        except FileNotFoundError:
            return {"success": False, "error": "文件不存在"}
        except PermissionError:
            return {"success": False, "error": "权限不足"}
        except Exception as e:
            return {"success": False, "error": str(e)}


class SandboxedFileOperationsTool:
    """沙箱文件操作工具"""
    
    def __init__(self, executor: "SandboxedToolExecutor"):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        operation = parameters.get("operation", "")
        path = parameters.get("path", "")
        content = parameters.get("content", "")
        
        # 路径安全检查
        if not self.executor._is_path_allowed(path):
            return {"success": False, "error": f"路径不在白名单中: {path}"}
        
        try:
            if operation == "read":
                with open(path, "r", encoding="utf-8") as f:
                    return {"success": True, "result": f.read()[:10000]}  # 限制读取大小
            
            elif operation == "write":
                # 检查文件大小
                if len(content) > 100000:
                    return {"success": False, "error": "内容超过限制"}
                
                with open(path, "w", encoding="utf-8") as f:
                    f.write(content)
                return {"success": True, "result": "文件写入成功"}
            
            elif operation == "append":
                if len(content) > 10000:
                    return {"success": False, "error": "内容超过限制"}
                
                with open(path, "a", encoding="utf-8") as f:
                    f.write(content)
                return {"success": True, "result": "内容追加成功"}
            
            elif operation == "list":
                import os
                files = os.listdir(path)
                return {"success": False, "result": files}
            
            else:
                return {"success": False, "error": f"未知操作: {operation}"}
        
        except FileNotFoundError:
            return {"success": False, "error": "文件不存在"}
        except PermissionError:
            return {"success": False, "error": "权限不足"}
        except Exception as e:
            return {"success": False, "error": str(e)}


class FeishuFileReadTool:
    """飞书文件读取工具 - 通过飞书API下载并读取文件内容"""
    
    def __init__(self, executor=None):
        self.executor = executor
    
    def _validate_message_id(self, message_id: str) -> bool:
        """验证 message_id 格式是否正确"""
        if not message_id:
            return True  # 允许为空（旧版文件可能不需要）
        
        # 飞书 message_id 格式：om_ + 32-64位字母数字字符（支持多种版本，包括包含x的格式）
        import re
        pattern = r"^om_[0-9a-z]{32,64}$"
        return bool(re.match(pattern, message_id.lower()))
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        file_key = parameters.get("file_key", "")
        user_id = parameters.get("user_id", "")
        message_id = parameters.get("message_id", "")
        
        # 记录详细请求参数
        logger.info(f"📋 收到文件读取请求 - file_key={file_key}, message_id={message_id}, user_id={user_id}")
        
        if not file_key:
            logger.error("❌ 文件读取失败：缺少file_key参数")
            return {"success": False, "error": "参数错误：缺少file_key参数"}
        
        # 验证 message_id 格式（如果提供了）
        if message_id and not self._validate_message_id(message_id):
            logger.error(f"❌ message_id 格式无效: {message_id}")
            logger.info(f"   正确格式示例: om_da5e1234567890abcdef1234567890ab")
            return {"success": False, "error": f"参数错误：message_id格式无效"}
        
        # 检查 file_v3 格式是否需要 message_id
        is_file_v3 = file_key.startswith("file_v3_")
        if is_file_v3 and not message_id:
            logger.error(f"❌ file_v3 格式文件需要有效的 message_id，当前 message_id 为空")
            return {"success": False, "error": "参数错误：file_v3格式文件需要message_id参数"}
        
        try:
            # 首先尝试从本地向量数据库获取文件内容（文件上传时已存储，支持用户隔离）
            file_content = self._get_content_from_vector_db(file_key, user_id)
            
            # 如果本地没有找到，再从飞书下载
            if not file_content:
                logger.info(f"本地存储未找到，尝试从飞书下载: {file_key}")
                
                # 优先使用executor的飞书客户端，如果没有则创建新的适配器
                api_client = None
                if self.executor and hasattr(self.executor, '_feishu_client'):
                    api_client = self.executor._feishu_client
                
                if api_client:
                    # 使用executor的客户端下载文件
                    file_content = self._download_feishu_file_direct(api_client, file_key, message_id)
                else:
                    from src.plugins.im_adapters import FeishuAdapter
                    # 创建飞书适配器
                    adapter = FeishuAdapter()
                    # 调用飞书API下载文件（使用推荐的 message resource 接口）
                    file_content = self._download_feishu_file(adapter, file_key, message_id)
            
            if file_content:
                # 将文件内容存储到向量数据库（可选）
                if user_id:
                    self._store_to_vector_db(user_id, file_content, file_key)
                
                logger.info(f"✅ 文件读取成功：file_key={file_key}, content_length={len(file_content)}")
                return {
                    "success": True,
                    "result": {
                        "file_key": file_key,
                        "content": file_content,
                        "content_length": len(file_content)
                    }
                }
            else:
                logger.error(f"❌ 文件下载失败：file_key={file_key}, message_id={message_id}")
                return {"success": False, "error": "文件下载失败，请检查飞书配置或网络连接"}
                
        except Exception as e:
            logger.error(f"❌ 读取飞书文件失败：file_key={file_key}, message_id={message_id}, 错误详情: {str(e)}", exc_info=True)
            return {"success": False, "error": f"文件读取失败: {str(e)}"}
    
    def _download_feishu_file_direct(self, api_client, file_key: str, message_id: str = None) -> Optional[str]:
        """
        直接使用api_client下载飞书文件（新版 file_v3 专用）
        支持：单聊、群聊、私聊、所有飞书文件格式
        
        Args:
            api_client: 飞书API客户端实例
            file_key: 文件标识（支持 file_v3_xxx 格式）
            message_id: 消息ID（用于调用推荐接口）
        
        Returns:
            文件内容字符串，失败返回None
        """
        try:
            # 检查 file_key 是否为 file_v3 格式
            is_file_v3 = file_key.startswith("file_v3_")
            
            # ======================
            # 方法 1：官方标准接口（唯一支持 file_v3）
            # ======================
            if message_id:
                try:
                    from lark_oapi.api.im.v1 import GetMessageResourceRequest, GetMessageResourceResponse
                    
                    req = GetMessageResourceRequest.builder() \
                        .message_id(message_id) \
                        .file_key(file_key) \
                        .type("file") \
                        .build()
                    
                    response: GetMessageResourceResponse = api_client.im.v1.message_resource.get(req)
                    
                    if response.success() and response.raw.content:
                        logger.info(f"📥 使用 message resource 接口下载文件成功: {file_key}")
                        logger.info(f"📦 文件内容长度: {len(response.raw.content)} bytes")
                        
                        # 尝试从响应头获取文件名
                        file_name = ""
                        if hasattr(response.raw, 'headers'):
                            content_disposition = response.raw.headers.get('Content-Disposition', '')
                            file_name = decode_filename(content_disposition)
                        logger.info(f"📝 提取文件名: '{file_name}'")
                        
                        # 解析文件内容
                        logger.info(f"🔍 开始解析文件内容...")
                        parsed_content = self._parse_file_content(response.raw.content, file_name)
                        logger.info(f"✅ 文件解析完成，内容长度: {len(parsed_content) if parsed_content else 0}")
                        return parsed_content
                    
                    logger.debug(f"message resource 接口未返回内容: code={response.code}, msg={response.msg}")
                except Exception as e:
                    logger.error(f"方法1(message resource)失败: {str(e)}")
            elif is_file_v3:
                # file_v3 格式必须有 message_id
                logger.error(f"file_v3 格式文件需要有效的 message_id，当前 message_id 为空")
                return None
            
            # ======================
            # 方法 2：云文档 API（兼容旧版，不支持 file_v3）
            # ======================
            if not is_file_v3:
                try:
                    from lark_oapi.api.drive.v1 import DownloadFileRequest
                    
                    req = DownloadFileRequest.builder().file_token(file_key).build()
                    response = api_client.drive.v1.file.download(req)
                    
                    if response.success() and response.raw.content:
                        logger.info(f"使用 drive API 下载文件成功: {file_key}")
                        
                        # 尝试从响应头获取文件名
                        file_name = ""
                        if hasattr(response.raw, 'headers'):
                            content_disposition = response.raw.headers.get('Content-Disposition', '')
                            file_name = decode_filename(content_disposition)
                        
                        # 解析文件内容
                        return self._parse_file_content(response.raw.content, file_name)
                    
                    logger.debug(f"drive API 未返回内容: code={response.code}, msg={response.msg}")
                except Exception as e:
                    logger.error(f"方法2(drive API)失败: {str(e)}")
            
            # 构建详细错误信息
            error_msg = f"飞书文件下载失败"
            if is_file_v3:
                error_msg += f"（file_v3格式需要有效的message_id）"
            else:
                error_msg += f"（请确认file_key正确）"
            
            raise Exception(error_msg)
            
        except Exception as e:
            logger.error(f"下载飞书文件失败：file_key={file_key}, message_id={message_id}, 错误详情: {str(e)}", exc_info=True)
            return None
    
    def _download_feishu_file(self, adapter, file_key: str, message_id: str = None) -> Optional[str]:
        """
        修复版：飞书新版 file_v3 专用下载方法（通过适配器）
        支持：单聊、群聊、私聊、所有飞书文件格式
        
        Args:
            adapter: 飞书适配器实例
            file_key: 文件标识（支持 file_v3_xxx 格式）
            message_id: 消息ID（用于调用推荐接口）
        
        Returns:
            文件内容字符串，失败返回None
        """
        try:
            # 检查API客户端是否已初始化
            if not hasattr(adapter, '_api_client') or adapter._api_client is None:
                logger.warning("飞书API客户端未初始化，尝试初始化")
                adapter._initialize_client()
            
            if adapter._api_client is None:
                logger.error(f"文件下载失败：无法初始化飞书API客户端，请检查.env文件中的FEISHU_APP_ID和FEISHU_APP_SECRET配置")
                return None
            
            # 检查 file_key 是否为 file_v3 格式
            is_file_v3 = file_key.startswith("file_v3_")
            
            # ======================
            # 方法 1：官方标准接口（唯一支持 file_v3）
            # ======================
            if message_id:
                try:
                    from lark_oapi.api.im.v1 import GetMessageResourceRequest, GetMessageResourceResponse
                    
                    req = GetMessageResourceRequest.builder() \
                        .message_id(message_id) \
                        .file_key(file_key) \
                        .type("file") \
                        .build()
                    
                    response: GetMessageResourceResponse = adapter._api_client.im.v1.message_resource.get(req)
                    
                    if response.success() and response.raw.content:
                        logger.info(f"使用 message resource 接口下载文件成功: {file_key}")
                        
                        # 尝试从响应头获取文件名
                        file_name = ""
                        if hasattr(response.raw, 'headers'):
                            content_disposition = response.raw.headers.get('Content-Disposition', '')
                            file_name = decode_filename(content_disposition)
                        
                        # 解析文件内容
                        return self._parse_file_content(response.raw.content, file_name)
                    
                    logger.debug(f"message resource 接口未返回内容: code={response.code}, msg={response.msg}")
                except Exception as e:
                    logger.error(f"方法1(message resource)失败: {str(e)}")
            elif is_file_v3:
                # file_v3 格式必须有 message_id
                logger.error(f"file_v3 格式文件需要有效的 message_id，当前 message_id 为空")
                return None
            
            # ======================
            # 方法 2：云文档 API（兼容旧版，不支持 file_v3）
            # ======================
            if not is_file_v3:
                try:
                    from lark_oapi.api.drive.v1 import DownloadFileRequest
                    
                    req = DownloadFileRequest.builder().file_token(file_key).build()
                    response = adapter._api_client.drive.v1.file.download(req)
                    
                    if response.success() and response.raw.content:
                        logger.info(f"使用 drive API 下载文件成功: {file_key}")
                        
                        # 尝试从响应头获取文件名
                        file_name = ""
                        if hasattr(response.raw, 'headers'):
                            content_disposition = response.raw.headers.get('Content-Disposition', '')
                            file_name = decode_filename(content_disposition)
                        
                        # 解析文件内容
                        return self._parse_file_content(response.raw.content, file_name)
                    
                    logger.debug(f"drive API 未返回内容: code={response.code}, msg={response.msg}")
                except Exception as e:
                    logger.error(f"方法2(drive API)失败: {str(e)}")
            
            # 构建详细错误信息
            error_msg = f"飞书文件下载失败"
            if is_file_v3:
                error_msg += f"（file_v3格式需要有效的message_id）"
            else:
                error_msg += f"（请确认file_key正确）"
            
            raise Exception(error_msg)
            
        except Exception as e:
            logger.error(f"下载飞书文件失败：file_key={file_key}, message_id={message_id}, 错误详情: {str(e)}", exc_info=True)
            return None
    
    def _parse_file_content(self, content: bytes, file_name: str) -> str:
        """
        根据文件类型解析内容
        """
        try:
            # 获取文件扩展名
            import os
            _, ext = os.path.splitext(file_name.lower())
            content_length = len(content)
            
            logger.debug(f"开始解析文件：file_name={file_name}, extension={ext}, content_length={content_length}")
            
            # 文本文件直接解码
            if ext in ['.txt', '.md', '.json', '.xml', '.html', '.csv', '.tsv']:
                logger.debug(f"解析文本文件：file_name={file_name}, encoding=utf-8")
                return content.decode('utf-8', errors='ignore')
            
            # docx 文件需要专门解析
            elif ext == '.docx':
                logger.debug(f"解析docx文件：file_name={file_name}")
                return self._parse_docx(content, file_name)
            
            # xlsx 文件解析
            elif ext == '.xlsx':
                logger.debug(f"解析xlsx文件：file_name={file_name}")
                return self._parse_xlsx(content, file_name)
            
            # xls 文件解析（旧版 Excel）
            elif ext == '.xls':
                logger.debug(f"解析xls文件：file_name={file_name}")
                return self._parse_xls(content, file_name)
            
            # pptx 文件解析
            elif ext == '.pptx':
                logger.debug(f"解析pptx文件：file_name={file_name}")
                return self._parse_pptx(content, file_name)
            
            # pdf 文件解析
            elif ext == '.pdf':
                logger.debug(f"解析PDF文件：file_name={file_name}")
                return self._parse_pdf(content, file_name)
            
            # rtf 文件解析
            elif ext == '.rtf':
                logger.debug(f"解析rtf文件：file_name={file_name}")
                return self._parse_rtf(content, file_name)
            
            # doc 文件解析（旧版 Word）
            elif ext == '.doc':
                logger.debug(f"解析doc文件：file_name={file_name}")
                return self._parse_doc(content, file_name)
            
            # 其他未知类型，尝试作为文本处理
            else:
                logger.warning(f"未知文件格式：file_name={file_name}, extension={ext}")
                # 先尝试 UTF-8 解码
                try:
                    text = content.decode('utf-8')
                    # 如果看起来像二进制数据，返回提示
                    if self._is_binary(text):
                        error_msg = f"无法解析文件格式: {ext}\n\n文件大小: {content_length} bytes\n\n支持的文件类型: .txt, .md, .json, .xml, .html, .csv, .tsv, .docx, .doc, .xlsx, .xls, .pptx, .pdf, .rtf"
                        logger.warning(f"文件为二进制格式：file_name={file_name}, extension={ext}")
                        return error_msg
                    return text
                except Exception as decode_e:
                    error_msg = f"无法解析文件格式: {ext}\n\n文件大小: {content_length} bytes\n\n支持的文件类型: .txt, .md, .json, .xml, .html, .csv, .tsv, .docx, .doc, .xlsx, .xls, .pptx, .pdf, .rtf"
                    logger.error(f"文件解码失败：file_name={file_name}, extension={ext}, 错误: {str(decode_e)}")
                    return error_msg
                    
        except Exception as e:
            error_msg = f"文件解析失败：{str(e)}\n\n文件名: {file_name}\n文件大小: {len(content)} bytes"
            logger.error(f"解析文件内容失败：file_name={file_name}, 错误详情: {str(e)}", exc_info=True)
            return error_msg
    
    def _parse_docx(self, content: bytes, file_name: str = "") -> str:
        """解析 docx 文件内容"""
        try:
            from docx import Document
            from io import BytesIO
            
            doc = Document(BytesIO(content))
            full_text = []
            
            for para in doc.paragraphs:
                full_text.append(para.text)
            
            parsed_content = '\n'.join(full_text)
            logger.debug(f"docx解析完成：file_name={file_name}, paragraphs={len(full_text)}, parsed_length={len(parsed_content)}")
            return parsed_content
            
        except Exception as e:
            error_msg = f"无法解析 docx 文件: {str(e)}\n\n文件名: {file_name}"
            logger.error(f"解析 docx 文件失败：file_name={file_name}, 错误详情: {str(e)}", exc_info=True)
            return error_msg
    
    def _parse_xlsx(self, content: bytes, file_name: str = "") -> str:
        """解析 xlsx 文件内容"""
        try:
            import pandas as pd
            from io import BytesIO
            
            df = pd.read_excel(BytesIO(content))
            parsed_content = df.to_string()
            logger.debug(f"xlsx解析完成：file_name={file_name}, rows={len(df)}, columns={len(df.columns)}, parsed_length={len(parsed_content)}")
            return parsed_content
            
        except Exception as e:
            error_msg = f"无法解析 xlsx 文件: {str(e)}\n\n文件名: {file_name}"
            logger.error(f"解析 xlsx 文件失败：file_name={file_name}, 错误详情: {str(e)}", exc_info=True)
            return error_msg
    
    def _parse_xls(self, content: bytes, file_name: str = "") -> str:
        """解析 xls 文件内容（旧版 Excel）"""
        try:
            import pandas as pd
            from io import BytesIO
            
            df = pd.read_excel(BytesIO(content), engine='xlrd')
            parsed_content = df.to_string()
            logger.debug(f"xls解析完成：file_name={file_name}, rows={len(df)}, columns={len(df.columns)}, parsed_length={len(parsed_content)}")
            return parsed_content
            
        except Exception as e:
            error_msg = f"无法解析 xls 文件: {str(e)}\n\n文件名: {file_name}"
            logger.error(f"解析 xls 文件失败：file_name={file_name}, 错误详情: {str(e)}", exc_info=True)
            return error_msg
    
    def _parse_pptx(self, content: bytes, file_name: str = "") -> str:
        """解析 pptx 文件内容"""
        try:
            from pptx import Presentation
            from io import BytesIO
            
            prs = Presentation(BytesIO(content))
            full_text = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        full_text.append(shape.text)
            
            parsed_content = '\n\n'.join(full_text)
            logger.debug(f"pptx解析完成：file_name={file_name}, slides={len(prs.slides)}, parsed_length={len(parsed_content)}")
            return parsed_content
            
        except Exception as e:
            error_msg = f"无法解析 pptx 文件: {str(e)}\n\n文件名: {file_name}"
            logger.error(f"解析 pptx 文件失败：file_name={file_name}, 错误详情: {str(e)}", exc_info=True)
            return error_msg
    
    def _parse_pdf(self, content: bytes, file_name: str = "") -> str:
        """解析 pdf 文件内容"""
        try:
            import pdfplumber
            from io import BytesIO
            
            logger.info(f"📄 开始解析PDF文件：file_name={file_name}, content_length={len(content)}")
            
            with pdfplumber.open(BytesIO(content)) as pdf:
                full_text = []
                page_count = len(pdf.pages)
                logger.info(f"📄 PDF页数: {page_count}")
                
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text:
                        full_text.append(text)
                        logger.debug(f"📄 第 {i+1} 页提取文本长度: {len(text)}")
            
            parsed_content = '\n\n'.join(full_text)
            logger.info(f"✅ PDF解析完成：file_name={file_name}, pages={page_count}, parsed_length={len(parsed_content)}")
            
            if not parsed_content:
                logger.warning(f"⚠️ PDF解析结果为空，可能是扫描件或加密PDF")
                return f"PDF文件内容为空或无法提取文本（可能是扫描件或加密PDF）\n\n文件名: {file_name}\n文件大小: {len(content)} bytes"
            
            return parsed_content
            
        except Exception as e:
            error_msg = f"无法解析 pdf 文件: {str(e)}\n\n文件名: {file_name}"
            logger.error(f"解析 pdf 文件失败：file_name={file_name}, 错误详情: {str(e)}", exc_info=True)
            return error_msg
    
    def _parse_rtf(self, content: bytes, file_name: str = "") -> str:
        """解析 rtf 文件内容"""
        try:
            # 简单的 RTF 文本提取
            text = content.decode('utf-8', errors='ignore')
            # 移除 RTF 控制字
            import re
            text = re.sub(r'\\[a-zA-Z]+[^{}]*', '', text)
            text = re.sub(r'\\[{}]', '', text)
            text = re.sub(r'{[^}]*}', '', text)
            parsed_content = text.strip()
            logger.debug(f"RTF解析完成：file_name={file_name}, parsed_length={len(parsed_content)}")
            return parsed_content
            
        except Exception as e:
            error_msg = f"无法解析 rtf 文件: {str(e)}\n\n文件名: {file_name}"
            logger.error(f"解析 rtf 文件失败：file_name={file_name}, 错误详情: {str(e)}", exc_info=True)
            return error_msg
    
    def _parse_doc(self, content: bytes, file_name: str = "") -> str:
        """解析 doc 文件内容（旧版 Word）"""
        try:
            # 使用 antiword 或其他方式，这里提供基础支持
            # 旧版 doc 是二进制格式，需要专门工具
            logger.warning(f"doc文件无法直接解析：file_name={file_name}, content_length={len(content)}")
            return f"无法直接解析 .doc 文件格式。请将文件另存为 .docx 格式后重试，或使用 Microsoft Word 打开复制内容。\n\n文件名: {file_name}\n文件大小: {len(content)} bytes"
            
        except Exception as e:
            error_msg = f"无法解析 doc 文件: {str(e)}\n\n文件名: {file_name}"
            logger.error(f"解析 doc 文件失败：file_name={file_name}, 错误详情: {str(e)}", exc_info=True)
            return error_msg
    
    def _is_binary(self, text: str) -> bool:
        """检测文本是否包含二进制数据"""
        # 检查是否有大量不可打印字符
        non_printable = sum(1 for c in text if ord(c) < 32 and c not in '\n\r\t')
        return non_printable > len(text) * 0.1
    
    def _store_to_vector_db(self, user_id: str, content: str, file_key: str):
        """将文件内容存储到向量数据库"""
        try:
            from src.data.database import db
            from src.types import MemoryEntry
            from src.utils import generate_id, get_timestamp
            
            # 存储到长期记忆
            memory_entry = MemoryEntry(
                id=generate_id(),
                user_id=user_id,
                type="long",
                content=content,
                timestamp=get_timestamp(),
                tags=["feishu_file", file_key]
            )
            db.save_memory(memory_entry)
            
            logger.info(f"文件内容已存储到向量数据库: {file_key}")
        except Exception as e:
            logger.warning(f"存储到向量数据库失败: {str(e)}")
    
    def _get_content_from_vector_db(self, file_key: str, user_id: str = None) -> Optional[str]:
        """
        从向量数据库中获取文件内容（支持用户隔离）
        
        Args:
            file_key: 文件标识
            user_id: 用户ID（可选，提供时只查询该用户的记录）
        
        Returns:
            文件内容字符串，如果未找到返回None
        """
        try:
            from src.data.database import db
            
            # 从记忆存储中查询包含该file_key的记录（支持用户隔离）
            memories = db.get_memories_by_tag(file_key, user_id)
            
            if memories:
                # 返回最新的内容
                memories.sort(key=lambda m: m.timestamp, reverse=True)
                logger.info(f"✅ 从本地存储获取文件内容成功: file_key={file_key}, user_id={user_id}")
                return memories[0].content
            
            logger.info(f"ℹ️ 未在向量数据库中找到文件: file_key={file_key}, user_id={user_id}")
            return None
            
        except Exception as e:
            logger.error(f"❌ 从向量数据库获取文件内容失败: {str(e)}")
            return None


# PPT生成工具类
class GenerateOutlineTool:
    """根据内容生成PPT大纲工具"""
    
    def __init__(self, executor=None):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        content = parameters.get("content", "")
        title = parameters.get("title", "文档总结")
        
        if not content:
            return {"success": False, "error": "参数错误：缺少content参数"}
        
        try:
            from src.plugins.model_routers import select_model, call_model
            
            model = select_model("summarization", "complex")
            if not model:
                model = select_model("document_analysis", "complex")
            
            if not model:
                return {"success": False, "error": "没有可用的模型来生成大纲"}
            
            prompt = f"""根据以下文档内容，生成一个结构化的PPT大纲：

文档内容：
{content}

要求：
1. 分析文档核心内容，提取关键章节（3-5个章节为宜）
2. 每个章节包含：章节标题 + 主要要点（3-5个）
3. 大纲结构清晰，逻辑连贯
4. 使用中文输出

PPT标题：{title}

请输出JSON格式，结构如下：
[
  {{"title": "章节1标题", "content": ["要点1", "要点2", "要点3"]}},
  {{"title": "章节2标题", "content": ["要点1", "要点2"]}}
]
"""
            result = call_model(model, [{"role": "user", "content": prompt}])
            
            import json
            try:
                outline = json.loads(result)
                if isinstance(outline, list) and len(outline) > 0:
                    logger.info(f"📋 大纲生成成功，共 {len(outline)} 个章节")
                    return {
                        "success": True,
                        "result": {
                            "outline": outline,
                            "title": title,
                            "chapters_count": len(outline)
                        }
                    }
                else:
                    return {"success": False, "error": f"无效的大纲格式: {result[:200]}"}
            except json.JSONDecodeError:
                return {"success": False, "error": f"无法解析大纲JSON: {result[:200]}"}
                
        except Exception as e:
            logger.error(f"大纲生成失败: {str(e)}", exc_info=True)
            return {"success": False, "error": f"大纲生成失败: {str(e)}"}


class GeneratePPTFromSlidesTool:
    """根据幻灯片列表生成PPT工具"""
    
    def __init__(self, executor=None):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        title = parameters.get("title", "Untitled Presentation")
        slides = parameters.get("slides", [])
        
        if not slides or not isinstance(slides, list):
            return {"success": False, "error": "参数错误：slides必须是非空数组"}
        
        try:
            from src.tools.ppt_generator import PPTGeneratorBase as PPTGenerator
            
            ppt_generator = PPTGenerator()
            output_path = ppt_generator.generate_ppt(title, slides)
            
            logger.info(f"📊 PPT生成成功: {output_path}")
            return {
                "success": True,
                "result": {
                    "file_path": output_path,
                    "title": title,
                    "slides_count": len(slides)
                }
            }
        except Exception as e:
            logger.error(f"PPT生成失败: {str(e)}", exc_info=True)
            return {"success": False, "error": f"PPT生成失败: {str(e)}"}


class GeneratePPTFromOutlineTool:
    """根据大纲生成PPT工具"""
    
    def __init__(self, executor=None):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        title = parameters.get("title", "Untitled Presentation")
        outline = parameters.get("outline", [])
        
        if not outline or not isinstance(outline, list):
            return {"success": False, "error": "参数错误：outline必须是非空数组"}
        
        try:
            from src.tools.ppt_generator import PPTGeneratorBase as PPTGenerator
            
            ppt_generator = PPTGenerator()
            output_path = ppt_generator.generate_from_outline(title, outline)
            
            logger.info(f"📊 PPT从大纲生成成功: {output_path}")
            return {
                "success": True,
                "result": {
                    "file_path": output_path,
                    "title": title,
                    "chapters_count": len(outline)
                }
            }
        except Exception as e:
            logger.error(f"PPT从大纲生成失败: {str(e)}", exc_info=True)
            return {"success": False, "error": f"PPT从大纲生成失败: {str(e)}"}


class GeneratePPTFromContentTool:
    """根据文件内容直接生成PPT工具（完整流程：内容→大纲→PPT）"""
    
    def __init__(self, executor=None):
        self.executor = executor
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        content = parameters.get("content", "")
        title = parameters.get("title", "文档总结")
        
        if not content:
            return {"success": False, "error": "参数错误：缺少content参数"}
        
        try:
            logger.info(f"📄 开始根据内容生成PPT，内容长度: {len(content)}")
            
            # 步骤1：根据内容生成大纲
            outline_tool = GenerateOutlineTool()
            outline_result = outline_tool.execute({"content": content, "title": title})
            
            if not outline_result.get("success"):
                return {"success": False, "error": f"大纲生成失败: {outline_result.get('error', '未知错误')}"}
            
            outline = outline_result["result"]["outline"]
            chapters_count = outline_result["result"]["chapters_count"]
            
            logger.info(f"📋 大纲生成成功，共 {chapters_count} 个章节")
            
            # 步骤2：根据大纲生成PPT
            from src.tools.ppt_generator import PPTGeneratorBase as PPTGenerator
            
            ppt_generator = PPTGenerator()
            output_path = ppt_generator.generate_from_outline(title, outline)
            
            logger.info(f"📊 PPT生成成功: {output_path}")
            return {
                "success": True,
                "result": {
                    "file_path": output_path,
                    "title": title,
                    "chapters_count": chapters_count,
                    "outline": outline
                }
            }
            
        except Exception as e:
            logger.error(f"根据内容生成PPT失败: {str(e)}", exc_info=True)
            return {"success": False, "error": f"根据内容生成PPT失败: {str(e)}"}


# 工具执行器注册表
TOOL_EXECUTOR_REGISTRY = {
    "basic": BasicToolExecutor,
    "sandboxed": SandboxedToolExecutor
}
