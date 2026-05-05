"""集成测试：验证PPT生成后自动发送文件"""

import sys
import os
import tempfile
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

from src.engine.ppt_workflow import ppt_workflow, WorkflowState
from src.gateway.message_router import MessageRouter
from src.types import Message
import json


def test_ppt_workflow_with_send():
    """测试PPT工作流完成后是否调用发送文件"""
    test_user_id = "test_ppt_send_integration"

    # 清理之前的上下文
    if hasattr(ppt_workflow, '_contexts') and test_user_id in ppt_workflow._contexts:
        del ppt_workflow._contexts[test_user_id]

    print("=== Step 1: 创建测试PPT文件 ===")
    # 创建一个真实的PPT文件
    from pptx import Presentation

    test_ppt_path = os.path.join(tempfile.gettempdir(), "test_ppt_send.pptx")
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(test_ppt_path)
    print(f"创建测试PPT: {test_ppt_path}")
    print(f"文件大小: {os.path.getsize(test_ppt_path)} bytes")

    print("\n=== Step 2: 模拟完成的工作流 ===")
    # 直接设置一个完成状态的上下文
    ctx = ppt_workflow._get_context(test_user_id)
    ctx.state = WorkflowState.COMPLETED
    ctx.output_path = test_ppt_path
    ctx.user_id = test_user_id
    ppt_workflow._update_mcp_context(test_user_id, ctx)

    print(f"Context state: {ctx.state}")
    print(f"Output path: {ctx.output_path}")

    print("\n=== Step 3: 测试 router._handle_ppt_generation ===")
    router = MessageRouter()

    # 检查 is_awaiting_confirmation
    is_waiting = ppt_workflow.is_awaiting_confirmation(test_user_id)
    print(f"is_awaiting_confirmation: {is_waiting}")

    # 注意：这里不能真正调用 router._handle_ppt_generation
    # 因为它需要完整的 IM 适配器

    # 但我们可以直接测试 _send_ppt_to_user_async 方法
    print("\n=== Step 4: 直接测试 _send_ppt_to_user_async ===")
    if hasattr(router, '_send_ppt_to_user_async'):
        print("调用 _send_ppt_to_user_async...")
        router._send_ppt_to_user_async(test_user_id, test_ppt_path)

        # 等待一下让线程执行
        import time
        time.sleep(2)

        print("方法调用完成，检查日志...")
    else:
        print("ERROR: _send_ppt_to_user_async 方法不存在")

    # 清理
    if os.path.exists(test_ppt_path):
        os.unlink(test_ppt_path)
        print(f"\n清理测试文件: {test_ppt_path}")

    print("\n=== 测试完成 ===")


if __name__ == "__main__":
    test_ppt_workflow_with_send()
